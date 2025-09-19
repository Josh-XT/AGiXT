from Extensions import Extensions
from agixtsdk import AGiXTSDK
from Globals import getenv
import os
import json
import requests
import subprocess
import sys
import asyncio
import logging
import cv2
import numpy as np
import tempfile
from soundfile import SoundFile, write as sf_write
try:
    from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, TextClip, CompositeVideoClip
except ImportError:
    import sys
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "moviepy"])
    from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, TextClip, CompositeVideoClip
try:
    from pptx import Presentation
except ImportError:
    import sys
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

class video(Extensions):
    """
    The Video extension for AGiXT. This extension provides capabilities for generating and editing videos
    from a list of images and narrations, using TTS for audio, MoviePy for video editing, and automatic frame
    duplication based on audio lengths for natural pacing.
    """
    CATEGORY = "Multimedia"
    
    def __init__(self, **kwargs):
        self.commands = {
            "Create Video from Images": self.create_video_from_images,
            "Add Text Overlay to Video": self.add_text_overlay,
            "Trim Video": self.trim_video,
            "Concatenate Videos": self.concatenate_videos,
            "Add Audio to Video": self.add_audio_to_video,
            "Adjust Video Speed": self.adjust_video_speed,
            "Apply Fade Transition": self.apply_fade_transition,
        }
        self.command_name = (
            kwargs["command_name"] if "command_name" in kwargs else "Video from Images"
        )
        self.user = kwargs["user"] if "user" in kwargs else ""
        self.agent_name = kwargs["agent_name"] if "agent_name" in kwargs else "gpt4free"
        self.conversation_name = (
            kwargs["conversation_name"] if "conversation_name" in kwargs else ""
        )
        self.WORKING_DIRECTORY = (
            kwargs["conversation_directory"]
            if "conversation_directory" in kwargs
            else os.path.join(os.getcwd(), "WORKSPACE")
        )
        os.makedirs(self.WORKING_DIRECTORY, exist_ok=True)
        self.conversation_id = (
            kwargs["conversation_id"] if "conversation_id" in kwargs else ""
        )
        self.ApiClient = (
            kwargs["ApiClient"]
            if "ApiClient" in kwargs
            else AGiXTSDK(
                base_uri=getenv("AGIXT_URI"),
                api_key=kwargs["api_key"] if "api_key" in kwargs else "",
            )
        )
        self.api_key = kwargs["api_key"] if "api_key" in kwargs else ""
        self.output_url = kwargs.get("output_url", "")
        self.failures = 0
    
    async def create_video_from_images(self, images_json: str, output_filename: str = "video.mp4", fps: int = 30, max_size_mb: int = 50):
        """
        Create a video from a list of images and narrations using MoviePy. Downloads images, generates TTS audio for each narration,
        and combines them into a video with automatic duration based on audio lengths.
        
        Args:
            images_json (str): JSON string of list of dicts: [{"image_url": "url1", "narration": "text1"}, ...]
            output_filename (str, optional): The name of the output video file. Defaults to "video.mp4".
            fps (int, optional): Frames per second for the video. Defaults to 30.
            max_size_mb (int, optional): Maximum size of the output video in MB. Defaults to 50.
        
        Returns:
            str: The URL of the generated video or error message.
        """
        try:
            # Parse JSON
            slides = json.loads(images_json)
            if not isinstance(slides, list) or len(slides) == 0:
                return "Error: Invalid images_json structure. Must be a list of dicts with 'image_url' and 'narration'."
            
            # Create temporary directory
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_dir_path = temp_dir
                os.makedirs(temp_dir_path, exist_ok=True)
                
                clips = []
                all_audio_data = []
                all_audio_lengths = []
                
                # Process images and narrations
                for idx, slide in enumerate(slides):
                    image_url = slide.get('image_url', '')
                    narration = slide.get('narration', '')
                    
                    if not image_url or not narration:
                        return f"Error: Missing 'image_url' or 'narration' for slide {idx + 1}."
                    
                    # Download image
                    img_path = os.path.join(temp_dir_path, f"image_{idx}.png")
                    img_response = requests.get(image_url)
                    img_response.raise_for_status()
                    with open(img_path, 'wb') as f:
                        f.write(img_response.content)
                    
                    # Generate TTS audio
                    audio_url = await self.ApiClient.text_to_speech(
                        text=narration,
                        agent_name=self.agent_name,
                    )
                    
                    # Download audio
                    audio_path = os.path.join(temp_dir_path, f"audio_{idx}.mp3")
                    audio_response = requests.get(audio_url)
                    audio_response.raise_for_status()
                    with open(audio_path, 'wb') as f:
                        f.write(audio_response.content)
                    
                    # Convert MP3 to WAV for duration calculation
                    wav_path = os.path.join(temp_dir_path, f"audio_{idx}.wav")
                    subprocess.run([
                        'ffmpeg', '-y', '-i', audio_path, '-ar', '22050', '-ac', '1', wav_path
                    ], capture_output=True, check=True)
                    
                    # Get audio duration
                    with SoundFile(wav_path) as audio_file:
                        audio_data = audio_file.read(dtype='float32')
                        sample_rate = audio_file.samplerate
                    
                    # Add 0.5s padding
                    padding = int(0.5 * sample_rate)
                    audio_data = np.pad(audio_data, (0, padding), mode='constant')
                    audio_duration = len(audio_data) / sample_rate
                    all_audio_data.append((audio_data, sample_rate))
                    all_audio_lengths.append(max(audio_duration, 2.0))
                    
                    # Create MoviePy ImageClip
                    clip = ImageClip(img_path, duration=max(audio_duration, 2.0))
                    audio_clip = AudioFileClip(audio_path)
                    clip = clip.set_audio(audio_clip)
                    clips.append(clip)
                
                # Concatenate clips
                output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
                final_clip = concatenate_videoclips(clips, method="compose")
                
                # Write video with initial quality
                final_clip.write_videofile(
                    output_path,
                    codec='libx264',
                    audio_codec='aac',
                    fps=fps,
                    preset='medium',
                    bitrate='2000k',
                )
                
                # Check size and compress if needed
                file_size_mb = os.path.getsize(output_path) / (1024 * 1024)
                if file_size_mb > max_size_mb:
                    final_clip.write_videofile(
                        output_path,
                        codec='libx264',
                        audio_codec='aac',
                        fps=fps,
                        preset='fast',
                        bitrate='1000k',
                    )
                    file_size_mb = os.path.getsize(output_path) / (1024 * 1024)
                    
                    if file_size_mb > max_size_mb:
                        new_fps = max(10, int(fps * (max_size_mb / file_size_mb) * 0.85))
                        final_clip.write_videofile(
                            output_path,
                            codec='libx264',
                            audio_codec='aac',
                            fps=new_fps,
                            preset='ultrafast',
                            bitrate='500k',
                        )
                
                final_clip.close()
                for clip in clips:
                    clip.close()
                
                return f"Video created successfully from {len(slides)} images. Size: {os.path.getsize(output_path)/(1024*1024):.2f}MB. Access at {self.output_url}{output_filename}"
        
        except Exception as e:
            return f"Error creating video from images: {str(e)}"
    
    async def add_text_overlay(self, video_path: str, text: str, output_filename: str = "text_overlay_video.mp4", position: str = "center", font_size: int = 50, font_color: str = "white"):
        """
        Add a text overlay to an existing video.
        
        Args:
            video_path (str): Path to the input video file in the workspace.
            text (str): Text to overlay on the video.
            output_filename (str, optional): Name of the output video file. Defaults to "text_overlay_video.mp4".
            position (str, optional): Position of the text ('center', 'top', 'bottom', 'left', 'right'). Defaults to "center".
            font_size (int, optional): Font size of the text. Defaults to 50.
            font_color (str, optional): Color of the text. Defaults to "white".
        
        Returns:
            str: The URL of the output video or error message.
        """
        try:
            input_path = os.path.join(self.WORKING_DIRECTORY, video_path)
            output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
            
            if not os.path.exists(input_path):
                return f"Error: Input video {video_path} does not exist."
            
            # Load video
            video_clip = VideoFileClip(input_path)
            
            # Create text clip
            text_clip = TextClip(
                text,
                fontsize=font_size,
                color=font_color,
                font='Arial',
                align=position,
                size=video_clip.size
            ).set_duration(video_clip.duration)
            
            # Composite video with text
            final_clip = CompositeVideoClip([video_clip, text_clip.set_position(position)])
            
            # Write output
            final_clip.write_videofile(
                output_path,
                codec='libx264',
                audio_codec='aac',
                fps=video_clip.fps
            )
            
            final_clip.close()
            video_clip.close()
            text_clip.close()
            
            return f"Text overlay added successfully. Access at {self.output_url}{output_filename}"
        
        except Exception as e:
            return f"Error adding text overlay: {str(e)}"
    
    async def trim_video(self, video_path: str, start_time: float, end_time: float, output_filename: str = "trimmed_video.mp4"):
        """
        Trim a video to a specified start and end time.
        
        Args:
            video_path (str): Path to the input video file in the workspace.
            start_time (float): Start time in seconds.
            end_time (float): End time in seconds.
            output_filename (str, optional): Name of the output video file. Defaults to "trimmed_video.mp4".
        
        Returns:
            str: The URL of the trimmed video or error message.
        """
        try:
            input_path = os.path.join(self.WORKING_DIRECTORY, video_path)
            output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
            
            if not os.path.exists(input_path):
                return f"Error: Input video {video_path} does not exist."
            
            # Load and trim video
            video_clip = VideoFileClip(input_path)
            if end_time > video_clip.duration:
                end_time = video_clip.duration
            trimmed_clip = video_clip.subclip(start_time, end_time)
            
            # Write output
            trimmed_clip.write_videofile(
                output_path,
                codec='libx264',
                audio_codec='aac',
                fps=video_clip.fps
            )
            
            trimmed_clip.close()
            video_clip.close()
            
            return f"Video trimmed successfully. Access at {self.output_url}{output_filename}"
        
        except Exception as e:
            return f"Error trimming video: {str(e)}"
    
    async def concatenate_videos(self, video_paths_json: str, output_filename: str = "concatenated_video.mp4"):
        """
        Concatenate multiple videos into one.
        
        Args:
            video_paths_json (str): JSON string of list of video file paths in the workspace: ["video1.mp4", "video2.mp4", ...]
            output_filename (str, optional): Name of the output video file. Defaults to "concatenated_video.mp4".
        
        Returns:
            str: The URL of the concatenated video or error message.
        """
        try:
            video_paths = json.loads(video_paths_json)
            if not isinstance(video_paths, list) or len(video_paths) == 0:
                return "Error: Invalid video_paths_json structure. Must be a list of video file paths."
            
            clips = []
            for video_path in video_paths:
                input_path = os.path.join(self.WORKING_DIRECTORY, video_path)
                if not os.path.exists(input_path):
                    return f"Error: Video {video_path} does not exist."
                clips.append(VideoFileClip(input_path))
            
            output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
            final_clip = concatenate_videoclips(clips, method="compose")
            
            final_clip.write_videofile(
                output_path,
                codec='libx264',
                audio_codec='aac',
                fps=clips[0].fps
            )
            
            final_clip.close()
            for clip in clips:
                clip.close()
            
            return f"Videos concatenated successfully. Access at {self.output_url}{output_filename}"
        
        except Exception as e:
            return f"Error concatenating videos: {str(e)}"
    
    async def add_audio_to_video(self, video_path: str, audio_text: str, output_filename: str = "video_with_audio.mp4"):
        """
        Add new audio narration to an existing video, replacing or overlaying existing audio.
        
        Args:
            video_path (str): Path to the input video file in the workspace.
            audio_text (str): Text for the new TTS audio.
            output_filename (str, optional): Name of the output video file. Defaults to "video_with_audio.mp4".
        
        Returns:
            str: The URL of the output video or error message.
        """
        try:
            input_path = os.path.join(self.WORKING_DIRECTORY, video_path)
            output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
            
            if not os.path.exists(input_path):
                return f"Error: Input video {video_path} does not exist."
            
            # Generate TTS audio
            audio_url = await self.ApiClient.text_to_speech(
                text=audio_text,
                agent_name=self.agent_name,
            )
            
            # Download audio
            with tempfile.TemporaryDirectory() as temp_dir:
                audio_path = os.path.join(temp_dir, "new_audio.mp3")
                audio_response = requests.get(audio_url)
                audio_response.raise_for_status()
                with open(audio_path, 'wb') as f:
                    f.write(audio_response.content)
                
                # Load video and audio
                video_clip = VideoFileClip(input_path)
                audio_clip = AudioFileClip(audio_path)
                
                # Adjust durations to match
                if audio_clip.duration > video_clip.duration:
                    audio_clip = audio_clip.subclip(0, video_clip.duration)
                else:
                    video_clip = video_clip.set_duration(audio_clip.duration)
                
                # Set new audio
                final_clip = video_clip.set_audio(audio_clip)
                
                # Write output
                final_clip.write_videofile(
                    output_path,
                    codec='libx264',
                    audio_codec='aac',
                    fps=video_clip.fps
                )
                
                final_clip.close()
                video_clip.close()
                audio_clip.close()
            
            return f"Audio added successfully. Access at {self.output_url}{output_filename}"
        
        except Exception as e:
            return f"Error adding audio to video: {str(e)}"
    
    async def adjust_video_speed(self, video_path: str, speed_factor: float, output_filename: str = "speed_adjusted_video.mp4"):
        """
        Adjust the playback speed of a video.
        
        Args:
            video_path (str): Path to the input video file in the workspace.
            speed_factor (float): Speed multiplier (e.g., 2.0 for double speed, 0.5 for half speed).
            output_filename (str, optional): Name of the output video file. Defaults to "speed_adjusted_video.mp4".
        
        Returns:
            str: The URL of the output video or error message.
        """
        try:
            input_path = os.path.join(self.WORKING_DIRECTORY, video_path)
            output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
            
            if not os.path.exists(input_path):
                return f"Error: Input video {video_path} does not exist."
            
            video_clip = VideoFileClip(input_path)
            final_clip = video_clip.speedx(speed_factor)
            
            final_clip.write_videofile(
                output_path,
                codec='libx264',
                audio_codec='aac',
                fps=video_clip.fps
            )
            
            final_clip.close()
            video_clip.close()
            
            return f"Video speed adjusted successfully. Access at {self.output_url}{output_filename}"
        
        except Exception as e:
            return f"Error adjusting video speed: {str(e)}"
    
    async def apply_fade_transition(self, video_path: str, transition_duration: float = 1.0, output_filename: str = "faded_video.mp4"):
        """
        Apply fade-in and fade-out transitions to a video.
        
        Args:
            video_path (str): Path to the input video file in the workspace.
            transition_duration (float, optional): Duration of fade transitions in seconds. Defaults to 1.0.
            output_filename (str, optional): Name of the output video file. Defaults to "faded_video.mp4".
        
        Returns:
            str: The URL of the output video or error message.
        """
        try:
            input_path = os.path.join(self.WORKING_DIRECTORY, video_path)
            output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
            
            if not os.path.exists(input_path):
                return f"Error: Input video {video_path} does not exist."
            
            video_clip = VideoFileClip(input_path)
            final_clip = video_clip.fadein(transition_duration).fadeout(transition_duration)
            
            final_clip.write_videofile(
                output_path,
                codec='libx264',
                audio_codec='aac',
                fps=video_clip.fps
            )
            
            final_clip.close()
            video_clip.close()
            
            return f"Fade transitions applied successfully. Access at {self.output_url}{output_filename}"
        
        except Exception as e:
            return f"Error applying fade transition: {str(e)}"

    async def extract_pptx_text(self, pptx_path: str) -> list[str]:
        """
        Extract text from each slide in a PPTX file, including slide content and notes.

        Args:
            pptx_path (str): Path to the PPTX file.

        Returns:
            list[str]: List of strings where each string is the concatenated text for that slide.
                       Returns an empty list if parsing fails.
        """
        try:
            presentation = Presentation(pptx_path)
            slides_text = []
            for slide in presentation.slides:
                slide_text = ""
                # Extract text from slide content
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                # Extract text from notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_text += slide.notes_slide.notes_text_frame.text + " "
                slides_text.append(slide_text.strip())
            return slides_text
        except Exception as e:
            return []
    async def extract_pptx_images(self, pptx_path: str, output_prefix: str = "photoname") -> list[str]:
        """
        Extract images from each slide in a PPTX file by converting slides to PNG images.

        Args:
            pptx_path (str): Path to the PPTX file.
            output_prefix (str, optional): Prefix for the output image filenames. Defaults to "photoname".

        Returns:
            list[str]: List of file paths to the saved PNG images. Returns an empty list if conversion fails.
        """
        try:
            presentation = Presentation(pptx_path)
            image_paths = []
            for idx, slide in enumerate(presentation.slides):
                image_path = os.path.join(self.WORKING_DIRECTORY, f"{output_prefix}{idx + 1}.png")
                await asyncio.to_thread(slide.export, image_path, 'PNG')
                image_paths.append(image_path)
            return image_paths
        except Exception as e:
            return []
    async def generate_slide_audios(self, texts: list[str], output_prefix: str = "audioname") -> list[str]:
        """
        Generate TTS audio for each slide text and download to workspace directory.

        Args:
            texts (list[str]): List of texts, one per slide.
            output_prefix (str, optional): Prefix for output filenames. Defaults to "audioname".

        Returns:
            list[str]: List of saved audio file paths.
        """
        audio_paths = []
        for idx, text in enumerate(texts):
            try:
                # Generate TTS audio
                audio_url = await self.ApiClient.text_to_speech(
                    text=text,
                    agent_name=self.agent_name,
                )
                # Download audio
                audio_path = os.path.join(self.WORKING_DIRECTORY, f"{output_prefix}{idx + 1}.mp3")
                audio_response = requests.get(audio_url)
                audio_response.raise_for_status()
                with open(audio_path, 'wb') as f:
                    f.write(audio_response.content)
                audio_paths.append(audio_path)
            except Exception as e:
                # Skip failed generations
                continue
        return audio_paths

    async def create_video_from_local_files(self, image_paths: list[str], audio_paths: list[str], output_filename: str = "Videoname.mp4", fps: int = 30) -> str:
        """
        Create a video from local image and audio files by pairing them by index.

        Args:
            image_paths (list[str]): List of paths to image files.
            audio_paths (list[str]): List of paths to audio files.
            output_filename (str, optional): Name of the output video file. Defaults to "Videoname.mp4".
            fps (int, optional): Frames per second for the video. Defaults to 30.

        Returns:
            str: Path to the created video file or error message.
        """
        try:
            min_len = min(len(image_paths), len(audio_paths))
            if len(image_paths) != len(audio_paths):
                logging.warning(f"Image and audio lists have different lengths: {len(image_paths)} vs {len(audio_paths)}. Using first {min_len} pairs.")

            clips = []
            for i in range(min_len):
                try:
                    img_clip = ImageClip(image_paths[i])
                    audio_clip = AudioFileClip(audio_paths[i])
                    img_clip = img_clip.set_duration(audio_clip.duration).set_audio(audio_clip)
                    clips.append(img_clip)
                except Exception as e:
                    logging.error(f"Error processing pair {i}: {e}")
                    continue

            if not clips:
                return "Error: No valid clips to create video."

            final_clip = concatenate_videoclips(clips)
            output_path = os.path.join(self.WORKING_DIRECTORY, output_filename)
            final_clip.write_videofile(output_path, fps=fps)

            final_clip.close()
            for clip in clips:
                clip.close()

            return output_path

        except Exception as e:
            return f"Error creating video from local files: {str(e)}"
    async def generate_video_from_pptx(self, pptx_path: str, output_filename: str = "Videoname.mp4") -> str:
        """
        Generate a video from a PPTX file by extracting text and images, generating audio for each slide,
        and combining them into a final video.

        Args:
            pptx_path (str): Path to the PPTX file.
            output_filename (str, optional): Name of the output video file. Defaults to "Videoname.mp4".

        Returns:
            str: Path to the generated video file or an error message.
        """
        try:
            # Extract text from PPTX
            texts = await self.extract_pptx_text(pptx_path)
            if not texts:
                return "Error: No text extracted from PPTX file."

            # Extract images from PPTX
            images = await self.extract_pptx_images(pptx_path)
            if not images:
                return "Error: No images extracted from PPTX file."

            # Generate audio for each text slide
            audios = await self.generate_slide_audios(texts)
            if not audios:
                return "Error: Failed to generate audio for any slides."

            # Check for mismatches
            if len(images) != len(audios):
                return f"Error: Mismatch in number of images ({len(images)}) and audios ({len(audios)})."

            # Create the final video
            video_path = await self.create_video_from_local_files(images, audios, output_filename)
            if video_path.startswith("Error"):
                return video_path

            return video_path

        except Exception as e:
            return f"Error generating video from PPTX: {str(e)}"