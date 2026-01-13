from Interactions import Interactions
from ApiClient import get_api_client, Conversations, Prompts, Chain
from Memories import Memories
from Extensions import Extensions
from Agent import get_agent_id_by_name
from pydub import AudioSegment
from Globals import getenv, get_tokens, DEFAULT_SETTINGS
from Models import ChatCompletions, TasksToDo, ChainCommandName, TranslationRequest
from Complexity import (
    calculate_complexity_score,
    ComplexityTier,
    ComplexityScore,
    log_complexity_decision,
)
from middleware import log_silenced_exception
from datetime import datetime
from typing import (
    List,
    Type,
    Union,
    get_args,
    get_origin,
    get_type_hints,
)
from MagicalAuth import MagicalAuth
from WorkerRegistry import worker_registry
from enum import Enum
from pydantic import BaseModel
from pptx import Presentation
from urllib.parse import urlparse
import ipaddress
import socket
import pdfplumber
import docx2txt
import zipfile
import pandas as pd
import subprocess
import logging
import asyncio
import requests
import inspect
import base64
import uuid
import json
import time
import os
import re


def is_safe_url(url: str) -> bool:
    """
    Validate a URL to prevent SSRF attacks.
    Blocks requests to internal networks, localhost, and cloud metadata endpoints,
    while allowing configured trusted local services (like ezlocalai).

    Args:
        url: The URL to validate

    Returns:
        bool: True if the URL is safe to request, False otherwise
    """
    try:
        parsed = urlparse(url)

        # Only allow http and https schemes
        if parsed.scheme not in ("http", "https"):
            logging.warning(f"SSRF protection: blocked non-http(s) scheme: {url}")
            return False

        hostname = parsed.hostname
        if not hostname:
            logging.warning(f"SSRF protection: no hostname in URL: {url}")
            return False

        port = parsed.port

        # Build list of trusted local URLs from environment configuration
        # These are internal services that AGiXT needs to communicate with
        trusted_local_urls = []

        # Check EZLOCALAI_URI / EZLOCALAI_API_URI for ezlocalai service
        ezlocalai_uri = getenv("EZLOCALAI_API_URI") or getenv("EZLOCALAI_URI")
        if ezlocalai_uri:
            trusted_local_urls.append(ezlocalai_uri)

        # Check AGIXT_URI for self-references
        agixt_uri = getenv("AGIXT_URI")
        if agixt_uri:
            trusted_local_urls.append(agixt_uri)

        # Check if the URL matches a trusted local service
        for trusted_url in trusted_local_urls:
            if trusted_url:
                try:
                    trusted_parsed = urlparse(trusted_url)
                    trusted_host = trusted_parsed.hostname
                    trusted_port = trusted_parsed.port

                    # Match if hostname and port match a trusted service
                    if hostname == trusted_host:
                        # If ports match (or trusted has no port and we're on default)
                        if port == trusted_port:
                            return True
                        # Also allow if trusted URL didn't specify a port
                        if trusted_port is None and port in (80, 443, None):
                            return True
                except Exception:
                    continue

        # Block cloud metadata endpoints (AWS, GCP, Azure, etc.)
        blocked_hosts = [
            "169.254.169.254",  # AWS/GCP metadata
            "metadata.google.internal",  # GCP metadata
            "metadata.google.com",
            "100.100.100.200",  # Alibaba Cloud metadata
            "169.254.170.2",  # AWS ECS task metadata
        ]
        if hostname in blocked_hosts:
            logging.warning(f"SSRF protection: blocked cloud metadata endpoint: {url}")
            return False

        # Resolve hostname to IP address(es)
        try:
            # Use getaddrinfo for both IPv4 and IPv6
            addr_info = socket.getaddrinfo(hostname, None, socket.AF_UNSPEC)
            ip_addresses = set()
            for info in addr_info:
                ip_str = info[4][0]
                ip_addresses.add(ip_str)
        except socket.gaierror:
            # DNS resolution failed - could be an invalid domain
            logging.warning(f"SSRF protection: DNS resolution failed for: {hostname}")
            return False

        # Check each resolved IP against blocked ranges
        for ip_str in ip_addresses:
            try:
                ip = ipaddress.ip_address(ip_str)

                # Block private networks
                if ip.is_private:
                    logging.warning(
                        f"SSRF protection: blocked private IP {ip_str} for URL: {url}"
                    )
                    return False

                # Block loopback addresses (127.0.0.0/8, ::1)
                if ip.is_loopback:
                    logging.warning(
                        f"SSRF protection: blocked loopback IP {ip_str} for URL: {url}"
                    )
                    return False

                # Block link-local addresses (169.254.0.0/16, fe80::/10)
                if ip.is_link_local:
                    logging.warning(
                        f"SSRF protection: blocked link-local IP {ip_str} for URL: {url}"
                    )
                    return False

                # Block multicast addresses
                if ip.is_multicast:
                    logging.warning(
                        f"SSRF protection: blocked multicast IP {ip_str} for URL: {url}"
                    )
                    return False

                # Block reserved addresses
                if ip.is_reserved:
                    logging.warning(
                        f"SSRF protection: blocked reserved IP {ip_str} for URL: {url}"
                    )
                    return False

            except ValueError:
                # Invalid IP address format
                logging.warning(f"SSRF protection: invalid IP address format: {ip_str}")
                return False

        return True

    except Exception as e:
        logging.error(f"SSRF protection: error validating URL {url}: {e}")
        return False


def sanitize_command_args_for_logging(command_args: dict) -> dict:
    """
    Sanitize command arguments for logging by redacting sensitive values.
    This prevents secrets from being exposed in conversation logs shown to users.

    Redacted keys (case-insensitive):
    - headers (may contain Authorization tokens)
    - Any key ending with _API_KEY
    - Any key ending with _SECRET
    - Authorization
    - Password
    - Any key containing 'secret' or 'password' or 'token' or 'key' (case-insensitive)

    Args:
        command_args: Dictionary of command arguments

    Returns:
        dict: Sanitized copy with sensitive values replaced by "[REDACTED]"
    """
    if not isinstance(command_args, dict):
        return command_args

    sensitive_patterns = [
        "headers",
        "authorization",
        "password",
        "passwd",
        "secret",
        "token",
        "credential",
        "private_key",
        "privatekey",
    ]

    sensitive_suffixes = [
        "_api_key",
        "_secret",
        "_token",
        "_password",
        "_key",
    ]

    sanitized = {}
    for key, value in command_args.items():
        key_lower = key.lower()

        # Check if key matches any sensitive pattern
        is_sensitive = False

        # Check exact/partial matches
        for pattern in sensitive_patterns:
            if pattern in key_lower:
                is_sensitive = True
                break

        # Check suffixes
        if not is_sensitive:
            for suffix in sensitive_suffixes:
                if key_lower.endswith(suffix):
                    is_sensitive = True
                    break

        if is_sensitive and value is not None and value != "":
            sanitized[key] = "[REDACTED]"
        elif isinstance(value, dict):
            # Recursively sanitize nested dicts
            sanitized[key] = sanitize_command_args_for_logging(value)
        else:
            sanitized[key] = value

    return sanitized


class AGiXT:
    def __init__(
        self,
        user: str,
        agent_name: str,
        api_key: str,
        conversation_name: str = None,
        collection_id=None,
    ):
        # Handle user dict from verify_api_key
        if isinstance(user, dict):
            user = user.get("email", "user")
        self.user_email = user.lower()
        if api_key is not None:
            api_key = str(api_key).replace("Bearer ", "").replace("bearer ", "")
        self.api_key = api_key
        self.auth = MagicalAuth(token=api_key)
        self.user_id = self.auth.user_id  # Cache user_id for reuse
        self.conversation = None
        self.conversation_id = None
        self.conversation_name = None
        if conversation_name != None:
            self.conversation = Conversations(
                conversation_name=conversation_name, user=self.user_email
            )
            self.conversation_id = self.conversation.get_conversation_id()
            self.conversation_name = conversation_name
        self.agent_name = agent_name
        self.uri = getenv("AGIXT_URI")
        if collection_id is not None:
            self.collection_id = str(collection_id)
        elif conversation_name:
            self.collection_id = self.conversation_id
        else:
            self.collection_id = "0"
        if self.conversation_name is None:
            self.conversation_name = datetime.now().strftime("%Y-%m-%d")
            self.conversation = Conversations(
                conversation_name=self.conversation_name, user=self.user_email
            )
            self.conversation_id = self.conversation.get_conversation_id()
        self.ApiClient = get_api_client(api_key)
        self.agent_interactions = Interactions(
            agent_name=self.agent_name,
            user=self.user_email,
            ApiClient=self.ApiClient,
            collection_id=self.collection_id,
        )
        self.agent = self.agent_interactions.agent
        self.agent_settings = (
            self.agent.AGENT_CONFIG["settings"]
            if "settings" in self.agent.AGENT_CONFIG
            else DEFAULT_SETTINGS
        )
        self.chain = Chain(user=self.user_email)
        self.prompts_manager = Prompts(user=self.user_email)  # Cache Prompts instance
        self.agent_workspace = self.agent.working_directory
        os.makedirs(self.agent_workspace, exist_ok=True)
        self.conversation_workspace = os.path.join(
            self.agent_workspace, self.conversation_id
        )
        os.makedirs(self.conversation_workspace, exist_ok=True)
        self.outputs = f"{self.uri}/outputs/{self.agent.agent_id}"
        self.failures = 0
        self.input_tokens = 0
        self.file_reader = None
        if self.collection_id is not None:
            self.file_reader = Memories(
                agent_name=self.agent_name,
                agent_config=self.agent.AGENT_CONFIG,
                collection_number=self.collection_id,
                ApiClient=self.ApiClient,
                user=self.user_email,
            )

    async def prompts(self, prompt_category: str = "Default"):
        """
        Get a list of available prompts

        Args:
            prompt_category (str): Category of the prompt

        Returns:
            list: List of available prompts
        """
        return self.prompts_manager.get_prompts(prompt_category=prompt_category)

    async def chains(self):
        """
        Get a list of available chains

        Returns:
            list: List of available chains
        """
        return self.chain.get_chains()

    async def settings(self):
        """
        Get the agent settings

        Returns:
            dict: Agent settings
        """
        return self.agent_settings

    async def commands(self):
        """
        Get a list of available commands

        Returns:
            list: List of available commands
        """
        return self.agent.available_commands

    async def browsed_links(self):
        """
        Get a list of browsed links

        Returns:
            list: List of browsed links
        """
        return self.agent.get_browsed_links()

    async def memories(
        self,
        user_input: str = "",
        limit_per_collection: int = 5,
        minimum_relevance_score: float = 0.3,
        additional_collection: str = "0",
    ):
        """
        Get a list of memories

        Args:
            user_input (str): User input to the agent
            limit_per_collection (int): Number of memories to return per collection
            minimum_relevance_score (float): Minimum relevance score for memories
            additional_collection (int): Additional collection number to pull memories from. Collections 0-5 are injected automatically.

        Returns:
            str: Agents relevant memories from the user input from collections 0-5 and the additional collection number if provided
        """
        formatted_prompt, prompt, tokens = await self.agent_interactions.format_prompt(
            user_input=user_input if user_input else "*",
            top_results=limit_per_collection,
            min_relevance_score=minimum_relevance_score,
            inject_memories_from_collection_number=additional_collection,
        )
        return formatted_prompt

    async def rename_new_conversation(self, user_input: str):
        """
        Rename a new conversation (one with name "-") based on the user input.
        This runs asynchronously and should be called right after user input is logged.
        The WebSocket will notify the frontend when the rename completes.

        Uses a direct LLM call (bypassing full inference pipeline) for speed.

        Args:
            user_input (str): The user's first message in the conversation
        """
        if self.conversation_name != "-":
            return  # Only rename new conversations

        try:
            # Use existing conversation instance instead of creating new one
            c = self.conversation

            # Default fallback name
            new_name = datetime.now().strftime("Conversation Created %Y-%m-%d %I:%M %p")

            # Get list of existing conversations to avoid duplicates
            conversation_list = c.get_conversations()

            # Build a simple, direct prompt for naming (bypasses Think About It pattern)
            naming_prompt = f"""Based on the user's message below, suggest a short, descriptive name for this conversation.

**User's message:**
{user_input[:500]}

**Existing conversation names to NOT use:**
{chr(10).join(conversation_list[:20])}

Respond with ONLY a JSON object in this exact format:
{{"suggested_conversation_name": "Your Suggested Name Here"}}

Rules:
- Use spaces in the name, not underscores
- Keep the name short (3-6 words)
- Make it descriptive of the topic
- Do not use any name from the existing list above
- Respond with ONLY the JSON, no explanation"""

            # Direct LLM call - bypasses full inference pipeline for speed
            new_convo = await self.agent.inference(
                prompt=naming_prompt,
                use_smartest=False,
                stream=False,
            )

            # Extract JSON from the response
            try:
                # Handle potential thinking tags in response (strip them out)
                if "<answer>" in new_convo:
                    new_convo = (
                        new_convo.split("<answer>")[-1].split("</answer>")[0].strip()
                    )

                if "```json" in new_convo:
                    json_text = new_convo.split("```json")[1].split("```")[0].strip()
                elif "```" in new_convo:
                    json_text = new_convo.split("```")[1].split("```")[0].strip()
                else:
                    json_start = new_convo.find("{")
                    json_end = new_convo.rfind("}")
                    if json_start != -1 and json_end != -1 and json_end > json_start:
                        json_text = new_convo[json_start : json_end + 1]
                    else:
                        raise ValueError("No valid JSON found in response")

                parsed_json = json.loads(json_text)
                new_name = parsed_json.get("suggested_conversation_name", new_name)

                # Handle duplicate names with a simpler retry
                if new_name in conversation_list:
                    retry_prompt = f"""The name "{new_name}" is already taken. Suggest a DIFFERENT name.

**User's message:**
{user_input[:300]}

Respond with ONLY: {{"suggested_conversation_name": "Different Name Here"}}"""

                    retry_response = await self.agent.inference(
                        prompt=retry_prompt,
                        use_smartest=False,
                        stream=False,
                    )

                    # Handle potential thinking tags
                    if "<answer>" in retry_response:
                        retry_response = (
                            retry_response.split("<answer>")[-1]
                            .split("</answer>")[0]
                            .strip()
                        )

                    if "```json" in retry_response:
                        json_text = (
                            retry_response.split("```json")[1].split("```")[0].strip()
                        )
                    elif "```" in retry_response:
                        json_text = (
                            retry_response.split("```")[1].split("```")[0].strip()
                        )
                    else:
                        json_start = retry_response.find("{")
                        json_end = retry_response.rfind("}")
                        if (
                            json_start != -1
                            and json_end != -1
                            and json_end > json_start
                        ):
                            json_text = retry_response[json_start : json_end + 1]
                        else:
                            raise ValueError("No valid JSON found in retry response")

                    parsed_json = json.loads(json_text)
                    new_name = parsed_json.get("suggested_conversation_name", new_name)

                    if new_name in conversation_list:
                        new_name = datetime.now().strftime(
                            "Conversation Created %Y-%m-%d %I:%M %p"
                        )

            except Exception as e:
                logging.error(f"Error parsing conversation name: {e}")
                if new_convo and isinstance(new_convo, str) and len(new_convo) < 100:
                    # Try to use as name if it looks reasonable
                    clean_name = new_convo.strip().strip('"').strip("'")
                    if len(clean_name) > 3 and len(clean_name) < 80:
                        new_name = clean_name

            # Apply the rename
            c.set_conversation_summary(summary=new_name)
            self.conversation_name = c.rename_conversation(new_name=new_name)

        except Exception as e:
            import traceback

            traceback.print_exc()
            logging.error(f"Error in rename_new_conversation: {e}")

    async def check_if_coding_required(self, user_input: str) -> bool:
        """
        Evaluates if coding is required to assist with the user's request.

        Args:
            user_input (str): The raw user input without context or instructions

        Returns:
            bool: True if coding is required, False otherwise
        """
        if not user_input:
            return False

        evaluation_prompt = f"""Analyze the following user request and determine if writing, modifying, debugging, executing code, or doing math is required to assist them.

User request: {user_input}

Respond with ONLY "true" if coding assistance is needed, or "false" if not.

Examples:
- "Can you help me debug this Python function?" -> true
- "Write a script to sort a list" -> true
- "Fix the bug in my code" -> true
- "Can you solve this math problem?" -> true
- "Can you count these items?" -> true
- "What's the weather like?" -> false
- "Explain how arrays work" -> false
- "Tell me about Python" -> false
- "Can you send an email to Joe?" -> false

Your response (true or false):"""

        try:
            # Use a lightweight inference call for evaluation
            response = await self.agent.PROVIDER.inference(prompt=evaluation_prompt)
            # Extract and normalize the response
            response = str(response).strip().lower()
            return "true" in response
        except Exception as e:
            logging.warning(f"Error checking if coding required: {str(e)}")
            # Default to False if we can't determine
            return False

    async def inference(
        self,
        user_input: str,
        prompt_category: str = "Default",
        prompt_name: str = "Custom Input",
        images: list = [],
        injected_memories: int = 100,
        conversation_results: int = 15,
        shots: int = 1,
        browse_links: bool = False,
        voice_response: bool = False,
        log_user_input: bool = True,
        log_output: bool = True,
        language: str = "en",
        command_overrides: list = None,
        **kwargs,
    ):
        """
        Run inference on the AGiXT agent

        Args:
            user_input (str): User input to the agent
            prompt_category (str): Category of the prompt
            prompt_name (str): Name of the prompt to use
            injected_memories (int): Number of memories to inject into the conversation
            conversation_results (int): Number of interactions to inject into the conversation
            browse_links (bool): Whether to browse links in the response
            images (list): List of image file paths
            shots (int): Number of responses to generate
            voice_response (bool): Whether to generate a voice response
            log_user_input (bool): Whether to log the user input
            log_output (bool): Whether to log the output
            language (str): Language of the response
            command_overrides (list): List of command overrides for the agent
            **kwargs: Additional keyword arguments

        Returns:
            str: Response from the agent
        """
        if "conversation_results" in kwargs:
            try:
                conversation_results = int(kwargs["conversation_results"])
            except:
                conversation_results = 10
            del kwargs["conversation_results"]
        if "context_results" in kwargs:
            try:
                injected_memories = int(kwargs["context_results"])
            except:
                injected_memories = 100
            del kwargs["context_results"]
        if "tts" in kwargs:
            voice_response = str(kwargs["tts"]).lower() == "true"
            del kwargs["tts"]
        if "conversation_name" in kwargs:
            del kwargs["conversation_name"]
        if "conversation_id" in kwargs:
            del kwargs["conversation_id"]

        # Calculate complexity score for inference-time compute scaling
        complexity_score = calculate_complexity_score(
            user_input=user_input,
            agent_settings=self.agent_settings,
        )

        # Log complexity decision for debugging
        # log_complexity_decision(complexity_score, user_input[:100] if user_input else "")

        # Determine use_smartest based on complexity scoring
        if "use_smartest" not in kwargs:
            kwargs["use_smartest"] = False
        if kwargs["use_smartest"] == False:
            # Use complexity-based routing instead of just coding check
            kwargs["use_smartest"] = complexity_score.route_to_smartest

        # Pass complexity score to interactions for thinking budget enforcement
        kwargs["complexity_score"] = complexity_score

        response = await self.agent_interactions.run(
            user_input=user_input,
            prompt_category=prompt_category,
            prompt_name=prompt_name,
            context_results=injected_memories,
            conversation_results=conversation_results,
            shots=shots,
            conversation_name=self.conversation_name,
            conversation_id=self.conversation_id,
            browse_links=browse_links,
            images=images,
            tts=voice_response,
            log_user_input=log_user_input,
            log_output=log_output,
            command_overrides=command_overrides,
            **kwargs,
        )
        if language == "en":
            return response
        translation_prompt = f"Markdown output is acceptable in the `target_language_translated_text` field, but all output should only be in the target language aside from that variable name. The goal is to intuitively and effectively translate the full given text from English to {language}. **Text to translate to target language '{language}'**:\n"
        target_language_user_input = await self.convert_to_model(
            input_string=f"{translation_prompt}{user_input}",
            model=TranslationRequest,
        )
        self.conversation.log_interaction(
            role="USER",
            message=f"{user_input}\n**Translated to {language.upper()}**\n{target_language_user_input.target_language_translated_text}",
        )
        if voice_response:
            await self.text_to_speech(
                text=target_language_user_input.target_language_translated_text,
                log_output=True,
            )
        stripped_response = re.sub(r"```[^```]+```", "", response)
        target_language_response = await self.convert_to_model(
            input_string=f"{translation_prompt}{stripped_response}",
            model=TranslationRequest,
        )
        new_response = f"{response}\n**Translated to {language.upper()}**\n{target_language_response.target_language_translated_text}"
        self.conversation.log_interaction(
            role=self.agent_name,
            message=new_response,
        )
        if voice_response:
            await self.text_to_speech(
                text=target_language_response.target_language_translated_text,
                log_output=True,
            )
        return new_response

    async def generate_image(self, prompt: str):
        """
        Generate an image from a prompt

        Args:
            prompt (str): Prompt for the image generation

        Returns:
            str: URL of the generated image
        """
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[ACTIVITY] Generating image.",
        )
        return await self.agent.generate_image(prompt=prompt)

    async def text_to_speech(
        self,
        text: str,
        log_output: bool = False,
    ):
        """
        Generate Text to Speech audio from text

        Args:
            text (str): Text to convert to speech
            log_output (bool): Whether to log the output

        Returns:
            str: URL of the generated audio
        """
        tts_url = await self.agent.text_to_speech(text=text)
        if not str(tts_url).startswith("http"):
            file_type = "wav"
            file_name = f"{uuid.uuid4().hex}.{file_type}"
            audio_path = os.path.join(self.agent_workspace, file_name)
            full_path = os.path.normpath(os.path.join(self.agent_workspace, file_name))
            if not full_path.startswith(self.agent_workspace):
                raise Exception("Path given not allowed")
            audio_data = base64.b64decode(tts_url)
            with open(audio_path, "wb") as f:
                f.write(audio_data)
            tts_url = f"{self.outputs}/{file_name}"
        if log_output:
            if tts_url.endswith(".mp3"):
                self.conversation.log_interaction(
                    role=self.agent_name,
                    message=f'<audio controls><source src="{tts_url}" type="audio/mpeg"></audio>',
                )
            else:
                self.conversation.log_interaction(
                    role=self.agent_name,
                    message=f'<audio controls><source src="{tts_url}" type="audio/wav"></audio>',
                )
        return tts_url

    async def audio_to_text(self, audio_path: str):
        """
        Audio to Text transcription

        Args:
            audio_path (str): Path to the audio file

        Returns
            str: Transcription of the audio
        """
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[ACTIVITY] Transcribing recorded audio.",
        )
        # Start a timer
        start = time.time()
        response = await self.agent.transcribe_audio(audio_path=audio_path)
        # End the timer
        end = time.time()
        elapsed_time = end - start
        elapsed_time = "{:.2f}".format(elapsed_time)
        self.conversation.update_message(
            message="[ACTIVITY] Transcribing recorded audio.",
            new_message=f"[ACTIVITY] Transcribed audio in {elapsed_time} seconds.",
        )
        return response

    async def translate_audio(self, audio_path: str):
        """
        Translate an audio file

        Args:
            audio_path (str): Path to the audio file

        Returns
            str: Translation of the audio
        """
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[ACTIVITY] Translating audio.",
        )
        response = await self.agent.translate_audio(audio_path=audio_path)
        return response

    async def execute_command(
        self,
        command_name: str,
        command_args: dict,
        voice_response: bool = False,
        log_output: bool = False,
        log_activities: bool = False,
    ):
        """
        Execute a command with arguments

        Args:
            command_name (str): Name of the command to execute
            command_args (dict): Arguments for the command
            voice_response (bool): Whether to generate a voice response
            log_output (bool): Whether to log the output
            log_activities (bool): Whether to log the activities

        Returns:
            str: Response from the command
        """
        if log_activities:
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[ACTIVITY] Executing command `{command_name}` with args:\n```json\n{json.dumps(sanitize_command_args_for_logging(command_args), indent=2)}```",
            )
            # Yield control to allow websocket handlers to send the update
            await asyncio.sleep(0)
        try:
            response = await Extensions(
                agent_name=self.agent_name,
                agent_id=self.agent.agent_id,
                agent_config=self.agent.AGENT_CONFIG,
                conversation_name=self.conversation_name,
                conversation_id=self.conversation_id,
                ApiClient=self.ApiClient,
                api_key=self.api_key,
                user=self.user_email,
            ).execute_command(
                command_name=command_name,
                command_args=command_args,
            )
        except Exception as e:
            logging.error(f"Error executing command: {e}")
            response = f"Error executing command `{command_name}`."
        if "tts_provider" in self.agent_settings and voice_response:
            if (
                self.agent_settings["tts_provider"] != "None"
                and self.agent_settings["tts_provider"] != ""
                and self.agent_settings["tts_provider"] != None
            ):
                await self.text_to_speech(
                    text=response,
                    log_output=log_output,
                )
        if log_output:
            self.conversation.log_interaction(role=self.agent_name, message=response)
        return response

    async def run_chain_step(
        self,
        chain_run_id=None,
        step: dict = {},
        chain_name="",
        user_input="",
        agent_override="",
        chain_args=None,
        running_command=None,
    ):
        if chain_args is None or not isinstance(chain_args, dict):
            chain_args = {}
        current_running_command = running_command
        if not current_running_command and isinstance(chain_args, dict):
            current_running_command = chain_args.get("running_command")
        if not current_running_command:
            current_running_command = chain_name
        if not chain_run_id:
            chain_run_id = await self.chain.get_chain_run_id(chain_name=chain_name)
        if step:
            if "prompt_type" in step:
                if agent_override != "":
                    agent_name = agent_override
                else:
                    agent_name = (
                        step["agent_name"] if "agent_name" in step else self.agent_name
                    )
                prompt_type = str(step["prompt_type"]).lower()
                step_number = step["step"]
                if "prompt_name" in step["prompt"]:
                    prompt_name = step["prompt"]["prompt_name"]
                else:
                    prompt_name = ""
                args = self.chain.get_step_content(
                    chain_run_id=chain_run_id,
                    chain_name=chain_name,
                    prompt_content=step["prompt"],
                    user_input=user_input,
                    agent_name=agent_name,
                )
                if chain_args:
                    for arg, value in chain_args.items():
                        # Only use chain_args value if:
                        # 1. The arg doesn't exist in args yet, OR
                        # 2. The arg exists but is empty/None AND chain_args has a real value, OR
                        # 3. The chain_args value is not empty/None (allowing explicit overrides)
                        # This prevents empty/None values from overwriting stored chain step arguments
                        existing_value = args.get(arg)
                        has_existing_value = (
                            existing_value is not None and existing_value != ""
                        )
                        has_new_value = value is not None and value != ""

                        if arg not in args:
                            # Arg doesn't exist, always add it
                            args[arg] = value
                        elif not has_existing_value and has_new_value:
                            # Existing is empty but new has value, use new
                            args[arg] = value
                        elif has_new_value:
                            # Both have values, new value overrides (explicit override)
                            args[arg] = value
                        # If existing has value and new is empty/None, keep existing (do nothing)
                log_output_flag = args.pop("log_output", None)
                if (
                    current_running_command
                    and not args.get("running_command")
                    and isinstance(args, dict)
                ):
                    args["running_command"] = current_running_command
                if "chain_name" in args:
                    args["chain"] = args["chain_name"]
                if "chain" not in args:
                    args["chain"] = chain_name
                if not args.get("conversation_name") and not args.get(
                    "conversation_id"
                ):
                    args["conversation_name"] = (
                        chain_args.get("conversation_name")
                        or self.conversation_name
                        or f"Chain Execution History: {chain_name}"
                    )
                if args.get("conversation"):
                    args["conversation_name"] = args["conversation"]
                # Ensure we have a conversation_id - it's more stable than name
                if not args.get("conversation_id"):
                    if chain_args.get("conversation_id"):
                        args["conversation_id"] = chain_args["conversation_id"]
                    elif self.conversation_id:
                        args["conversation_id"] = self.conversation_id
                    elif args.get("conversation_name"):
                        # Get the conversation_id from the name if we only have a name
                        from Conversations import Conversations as ConvHelper

                        conv_helper = ConvHelper(
                            conversation_name=args["conversation_name"], user=self.user
                        )
                        args["conversation_id"] = conv_helper.get_conversation_id()
                if prompt_type == "command":
                    self.conversation.log_interaction(
                        role=self.agent_name,
                        message=f"[SUBACTIVITY] Executing command `{step['prompt']['command_name']}` with args:\n```json\n{json.dumps(sanitize_command_args_for_logging(args), indent=2)}```",
                    )
                    # Yield control to allow websocket handlers to send the update
                    await asyncio.sleep(0)
                    result = await self.execute_command(
                        command_name=step["prompt"]["command_name"],
                        command_args=args,
                        voice_response=False,
                    )
                elif prompt_type == "prompt":
                    if "command_name" in args:
                        del args["command_name"]
                    if "chain_name" in args:
                        del args["chain_name"]
                    if "chain_args" in args:
                        del args["chain_args"]
                    if "chain" in args:
                        del args["chain"]
                    # Always use conversation_id as conversation_name since IDs are stable and unique
                    args["conversation_name"] = args.get("conversation_id") or args.get(
                        "conversation_name"
                    )
                    if prompt_name == "":
                        prompt_name = "Think About It"
                    prompt_args = args.copy()
                    if "browse_links" not in prompt_args:
                        prompt_args["browse_links"] = False
                    if current_running_command:
                        prompt_args["running_command"] = current_running_command
                    prompt_args["prompt_name"] = prompt_name
                    prompt_args["log_user_input"] = False
                    prompt_args["voice_response"] = False
                    if log_output_flag is None:
                        prompt_args["log_output"] = False
                    else:
                        prompt_args["log_output"] = str(
                            log_output_flag
                        ).lower() not in [
                            "false",
                            "0",
                            "no",
                        ]
                    prompt_args["user_input"] = user_input
                    self.conversation.log_interaction(
                        role=self.agent_name,
                        message=f"[SUBACTIVITY] Running prompt: `{prompt_name}` with args:\n```json\n{json.dumps(prompt_args, indent=2)}```",
                    )
                    # Yield control to allow websocket handlers to send the update
                    await asyncio.sleep(0)
                    # Get agent_id from agent_name for the API call
                    step_agent_id = get_agent_id_by_name(
                        agent_name=agent_name, user=self.user_email
                    )
                    result = self.ApiClient.prompt_agent(
                        agent_id=step_agent_id,
                        prompt_name=prompt_name,
                        prompt_args=prompt_args,
                    )
                elif prompt_type == "chain":
                    self.conversation.log_interaction(
                        role=self.agent_name,
                        message=f"[SUBACTIVITY] Running chain: `{args['chain']}` with args:\n```json\n{json.dumps(sanitize_command_args_for_logging(args), indent=2)}```",
                    )
                    # Yield control to allow websocket handlers to send the update
                    await asyncio.sleep(0)
                    if "chain_name" in args:
                        args["chain"] = args["chain_name"]
                    if "user_input" in args:
                        args["input"] = args["user_input"]
                    elif "input" not in args:
                        args["input"] = user_input
                    if isinstance(args.get("chain_args"), dict):
                        nested_chain_args = args["chain_args"].copy()
                    elif "conversation_name" in args:
                        nested_chain_args = {
                            "conversation_name": args["conversation_name"]
                        }
                    else:
                        nested_chain_args = {}
                    if (
                        current_running_command
                        and "running_command" not in nested_chain_args
                    ):
                        nested_chain_args["running_command"] = current_running_command
                    if (
                        log_output_flag is not None
                        and "log_output" not in nested_chain_args
                    ):
                        nested_chain_args["log_output"] = str(
                            log_output_flag
                        ).lower() not in [
                            "false",
                            "0",
                            "no",
                        ]
                    result = await self.execute_chain(
                        chain_name=args["chain"],
                        user_input=args["input"],
                        agent_override=agent_name,
                        from_step=args["from_step"] if "from_step" in args else 1,
                        chain_args=nested_chain_args,
                        running_command=current_running_command,
                        log_user_input=False,
                        voice_response=False,
                    )
        if result:
            if isinstance(result, dict) and "response" in result:
                result = result["response"]
            if result == "Unable to retrieve data.":
                result = None
            if isinstance(result, dict):
                result = json.dumps(result)
            if not isinstance(result, str):
                result = str(result)
            await self.chain.update_step_response(
                chain_run_id=chain_run_id,
                chain_name=chain_name,
                step_number=step_number,
                response=result,
            )
            return result
        else:
            return None

    async def execute_chain(
        self,
        chain_name,
        chain_run_id=None,
        user_input=None,
        agent_override="",
        from_step=1,
        chain_args=None,
        running_command=None,
        log_user_input=False,
        log_output=True,
        voice_response=False,
    ):
        if isinstance(chain_args, dict):
            merged_chain_args = chain_args.copy()
        else:
            merged_chain_args = {}
        if self.conversation_name and not merged_chain_args.get("conversation_name"):
            merged_chain_args["conversation_name"] = self.conversation_name
        if self.conversation_id and not merged_chain_args.get("conversation_id"):
            merged_chain_args["conversation_id"] = self.conversation_id
        active_running_command = (
            running_command or merged_chain_args.get("running_command") or chain_name
        )
        if active_running_command:
            merged_chain_args["running_command"] = active_running_command
        chain_data = self.chain.get_chain(chain_name=chain_name)
        if not chain_run_id:
            chain_run_id = await self.chain.get_chain_run_id(chain_name=chain_name)
        if chain_data == {}:
            return f"Chain `{chain_name}` not found."
        if log_user_input:
            self.conversation.log_interaction(
                role="USER",
                message=user_input,
            )
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[SUBACTIVITY] Running chain `{chain_name}`.",
        )
        # Yield control to allow websocket handlers to send the update
        await asyncio.sleep(0)
        response = ""
        step_responses = []
        step_summaries = []
        if "steps" not in chain_data:
            return f"Chain `{chain_name}` has no steps."
        if len(chain_data["steps"]) == 0:
            return f"Chain `{chain_name}` has no steps."
        for step_data in chain_data["steps"]:
            if int(step_data["step"]) >= int(from_step):
                if "prompt" in step_data and "step" in step_data:
                    step = {}
                    if "agent_name" not in step_data:
                        step_data["agent_name"] = self.agent_name
                    step["agent_name"] = (
                        agent_override
                        if agent_override != ""
                        else step_data["agent_name"]
                    )
                    step["prompt_type"] = step_data["prompt_type"]
                    step["prompt"] = step_data["prompt"]
                    step["step"] = step_data["step"]
                    task = await self.run_chain_step(
                        chain_run_id=chain_run_id,
                        step=step,
                        chain_name=chain_name,
                        user_input=user_input,
                        agent_override=agent_override,
                        chain_args=merged_chain_args,
                        running_command=active_running_command,
                    )
                    step_responses.append(task)
                    if task:
                        if not isinstance(task, str):
                            try:
                                task = json.dumps(task)
                            except Exception:
                                task = str(task)
                        step_type = step.get("prompt_type", "").title()
                        step_prompt = step.get("prompt", {})
                        step_identifier = ""
                        if isinstance(step_prompt, dict):
                            step_identifier = (
                                step_prompt.get("prompt_name")
                                or step_prompt.get("command_name")
                                or step_prompt.get("chain_name")
                                or ""
                            )
                        if step_identifier:
                            step_label = (
                                f"Step {step['step']} ({step_type} - {step_identifier})"
                            )
                        else:
                            step_label = f"Step {step['step']} ({step_type})"
                        step_summaries.append(f"{step_label} Output:\n{task}")
        if step_responses:
            response = step_responses[-1]
            # Only include step labels/summaries if there are multiple steps
            # For single-step chains (like commands), just return the clean output
            if step_summaries and len(step_summaries) > 1:
                response = "\n\n".join(step_summaries)
        if response == None:
            return f"Chain failed to complete, it failed on step {step_data['step']}. You can resume by starting the chain from the step that failed with chain ID {chain_run_id}."
        if log_output:
            self.conversation.log_interaction(role=self.agent_name, message=response)
        if "tts_provider" in self.agent_settings and voice_response:
            if (
                self.agent_settings["tts_provider"] != "None"
                and self.agent_settings["tts_provider"] != ""
                and self.agent_settings["tts_provider"] != None
            ):
                await self.text_to_speech(text=response, log_output=True)
        return response

    async def learn_from_websites(
        self,
        urls: list = [],
        summarize_content: bool = False,
    ):
        """
        Scrape a website and summarize the content

        Args:
            urls (list): List of URLs to scrape
            scrape_depth (int): Depth to scrape each URL
            summarize_content (bool): Whether to summarize the content

        Returns:
            str: Agent response with a list of scraped links
        """
        if isinstance(urls, str):
            user_input = f"Learn from the information from this website:\n {urls} "
        else:
            url_str = {"\n".join(urls)}
            user_input = f"Learn from the information from these websites:\n {url_str} "
        response = await self.agent_interactions.websearch.scrape_websites(
            user_input=user_input,
            summarize_content=summarize_content,
            conversation_name=self.conversation_name,
        )
        return (
            "I have scraped the information from the websites and saved it to memory."
        )

    def _get_large_file_instructions(
        self,
        file_name: str,
        file_path: str,
        file_size_bytes: int,
        converted_file_name: str = None,
        file_type: str = None,
    ) -> str:
        """
        Generate instructions for the LLM to read a large file that exceeds the token limit.

        Args:
            file_name: Original uploaded file name
            file_path: Path to the file in workspace
            file_size_bytes: Size of the file in bytes
            converted_file_name: Name of converted file if applicable (e.g., CSV from Excel)
            file_type: Type of file for specific instructions

        Returns:
            str: Instructions for the LLM on how to work with the file
        """
        file_size_mb = round(file_size_bytes / (1024 * 1024), 2)
        actual_file = converted_file_name if converted_file_name else file_name

        # Build the instruction based on file type
        if file_type and file_type.lower() in ["csv", "xlsx", "xls"]:
            return (
                f"The user uploaded a file called `{file_name}` which is {file_size_mb}MB in size. "
                f"The file has been converted to CSV format and saved as `{actual_file}` in the workspace. "
                f"Since this file is too large to include in context, use the **Search Files** command "
                f"to search for specific data patterns or column names, or use the **Read File Lines** command "
                f"to read specific line ranges from `{actual_file}`. Start by reading the first 50 lines "
                f"to understand the file structure and column headers."
            )
        elif file_type and file_type.lower() in [
            "md",
            "txt",
            "json",
            "py",
            "js",
            "ts",
            "html",
            "css",
            "yaml",
            "yml",
        ]:
            return (
                f"The user uploaded a file called `{file_name}` which is {file_size_mb}MB in size. "
                f"Since this file is too large to include in context, use the **Search Files** command "
                f"to search for specific content, or use the **Read File Lines** command to read specific "
                f"line ranges from `{actual_file}`. The file is a text-based file that can be read directly."
            )
        elif file_type and file_type.lower() == "pdf":
            return (
                f"The user uploaded a PDF file called `{file_name}` which is {file_size_mb}MB in size. "
                f"The PDF content has been extracted and saved to the workspace. Since the full content is too large to include "
                f"in context, use the **Search Files** command to search for specific content, or use the **Read File Lines** command "
                f"to read specific line ranges from the file."
            )
        else:
            return (
                f"The user uploaded a file called `{file_name}` which is {file_size_mb}MB in size. "
                f"Since this file is too large to include in context, use the **Search Files** command "
                f"to search for specific content, or use the **Read File Lines** command to read specific "
                f"line ranges from the file."
            )

    def _get_file_access_instructions(
        self,
        file_name: str,
        file_url: str,
        file_path: str,
        file_size_bytes: int,
        file_tokens: int,
        converted_file_name: str = None,
        file_type: str = None,
    ) -> str:
        """
        Generate instructions for the LLM to access an uploaded file using commands.
        This is used instead of including file content directly in context.

        Args:
            file_name: Original uploaded file name
            file_url: URL where the file can be accessed
            file_path: Path to the file in workspace
            file_size_bytes: Size of the file in bytes
            file_tokens: Number of tokens the file content would take
            converted_file_name: Name of converted file if applicable (e.g., CSV from Excel)
            file_type: Type of file for specific instructions

        Returns:
            str: Metadata and instructions for accessing the file
        """
        file_size_kb = round(file_size_bytes / 1024, 1)

        # Get the relative path from the conversation directory for use in commands
        # The file_path is an absolute path like: agent_workspace/conversation_id/filename
        # Commands run with WORKING_DIRECTORY = agent_workspace/conversation_id
        # So we need just the filename, not the conversation_id prefix
        conversation_dir = os.path.join(self.agent_workspace, self.conversation_id)
        if file_path.startswith(conversation_dir):
            # File is in the conversation directory - use path relative to it
            relative_path = file_path[len(conversation_dir) :].lstrip("/\\")
        elif file_path.startswith(self.agent_workspace):
            # File is in agent workspace but not conversation dir - use path relative to agent workspace
            relative_path = file_path[len(self.agent_workspace) :].lstrip("/\\")
        else:
            # Fallback to just the filename
            relative_path = file_name

        # Use converted file name if available, but with correct path
        if converted_file_name and converted_file_name != file_name:
            # Replace the filename in the relative path with the converted filename
            dir_path = os.path.dirname(relative_path)
            actual_path = (
                os.path.join(dir_path, converted_file_name)
                if dir_path
                else converted_file_name
            )
        else:
            actual_path = relative_path

        # If actual_path is empty, fall back to filename
        if not actual_path:
            actual_path = converted_file_name if converted_file_name else file_name

        # Base info about the file
        # For xlsx/xls files that were converted, emphasize the CSV as the primary file to use
        if (
            converted_file_name
            and converted_file_name != file_name
            and file_type
            and file_type.lower() in ["xlsx", "xls"]
        ):
            info = f"## File Available: `{converted_file_name}` (converted from `{file_name}`)\n"
            info += f"- **Original file:** `{file_name}` (Excel format - do NOT read this directly)\n"
            info += f"- **Use this file:** `{actual_path}` (CSV format - use this for all commands)\n"
            info += f"- **Size:** {file_size_kb} KB ({file_tokens} tokens)\n"
            info += f"- **URL:** [{file_name}]({file_url})\n"
            info += "\n**IMPORTANT:** The Excel file has been converted to CSV. Always use the CSV file (`{actual_path}`) for Read File and pandas operations.\n"
        else:
            info = f"## Uploaded File: `{file_name}`\n"
            info += f"- **Size:** {file_size_kb} KB ({file_tokens} tokens)\n"
            info += f"- **Path for commands:** `{actual_path}`\n"
            info += f"- **URL:** [{file_name}]({file_url})\n"

            if converted_file_name and converted_file_name != file_name:
                info += f"- **Converted to:** `{converted_file_name}` (CSV format)\n"

        info += "\n"

        # Type-specific instructions
        if file_type and file_type.lower() in ["csv", "xlsx", "xls"]:
            info += "### How to Access This Data:\n"
            info += f"1. **Read the file** using `Read File` command with filename `{actual_path}` to see the full content\n"
            info += "2. **Analyze with Python** using `Execute Python Code` to load with pandas and perform analysis:\n"
            info += f"   ```python\n   import pandas as pd\n   df = pd.read_csv('{actual_path}')\n   print(df.head())\n   print(df.describe())\n   ```\n"
            info += "3. **Search for specific data** using `Search File Content` to find particular values\n"
        elif file_type and file_type.lower() in [
            "py",
            "js",
            "ts",
            "json",
            "yaml",
            "yml",
            "md",
            "txt",
            "html",
            "css",
        ]:
            info += "### How to Access This File:\n"
            info += f"1. **Read the file** using `Read File` command with filename `{actual_path}`\n"
            info += "2. **Search for patterns** using `Search File Content` to find specific code or text\n"
            if file_type.lower() == "py":
                info += f"3. **Execute the code** using `Execute Python File` with path `{actual_path}`\n"
        elif file_type and file_type.lower() == "pdf":
            info += "### How to Access This PDF:\n"
            info += "The PDF content has been extracted and is available in the workspace.\n"
            info += f"1. **Read the file** using `Read File` command with filename `{actual_path}`\n"
            info += "2. **Search for patterns** using `Search File Content` to find specific text\n"
        else:
            info += "### How to Access This File:\n"
            info += f"1. **Read the file** using `Read File` command with filename `{actual_path}`\n"
            info += "2. **Search for content** using `Search File Content`\n"

        return info

    async def learn_spreadsheet(
        self, user_input, file_path, thinking_id, save_to_memory: bool = False
    ):
        file_name = os.path.basename(file_path)
        file_type = str(file_name).split(".")[-1]
        string_file_content = ""
        action_verb = "Learned" if save_to_memory else "Saved"
        action_location = "to memory" if save_to_memory else "to workspace"
        try:
            if file_type.lower() == "csv":
                df = pd.read_csv(file_path)
                csv = df.to_csv(index=False)
                string_file_content += f"Content from file uploaded named `{file_name}`:\n```csv\n{csv}```\n"
                return (
                    f"{action_verb} [{file_name}]({file_path}) {action_location}.",
                    string_file_content,
                )
            else:  # Excel file
                try:
                    xl = pd.ExcelFile(file_path)
                    if len(xl.sheet_names) > 1:
                        sheet_count = len(xl.sheet_names)
                        for i, sheet_name in enumerate(xl.sheet_names, 1):
                            df = xl.parse(sheet_name)
                            csv_file_path = file_path.replace(
                                f".{file_type}", f"_{i}.csv"
                            )
                            csv_file_name = os.path.basename(csv_file_path)
                            df.to_csv(csv_file_path, index=False)
                            message, file_content = await self.learn_spreadsheet(
                                user_input=user_input,
                                file_path=csv_file_path,
                                thinking_id=thinking_id,
                                save_to_memory=save_to_memory,
                            )
                            string_file_content += file_content
                        return (
                            f"Processed all sheets in [{file_name}]({file_path}).",
                            string_file_content,
                        )
                    else:
                        # Single sheet - also save as CSV for easier access
                        df = pd.read_excel(file_path)
                        csv = df.to_csv(index=False)
                        csv_file_path = file_path.replace(f".{file_type}", ".csv")
                        csv_file_name = os.path.basename(csv_file_path)
                        df.to_csv(csv_file_path, index=False)
                        string_file_content += f"Content from uploaded Excel file `{file_name}` (converted and saved as `{csv_file_name}` - use this CSV file for all Read File and pandas operations):\n```csv\n{csv}```\n"
                        return (
                            f"Converted [{file_name}]({file_path}) to CSV format at [{csv_file_name}]({csv_file_path}). Use `{csv_file_name}` for file operations.",
                            string_file_content,
                        )
                except Exception as e:
                    self.conversation.log_interaction(
                        role=self.agent_name,
                        message=f"[ACTIVITY][ERROR] Failed to read Excel file `{file_name}`: {str(e)}",
                    )
                    return (
                        f"Failed to read [{file_name}]({file_path}). Error: {str(e)}",
                        "",
                    )
        except Exception as e:
            logging.error(f"Unexpected error processing spreadsheet: {e}")
            return f"Failed to process [{file_name}]({file_path}). Unexpected error: {str(e)}"

    async def learn_from_file(
        self,
        file_url: str = "",
        file_name: str = "",
        user_input: str = "",
        collection_id: str = "0",
        thinking_id: str = "",
        save_to_memory: bool = False,
    ):
        """
        Learn from a file

        Args:
            file_url (str): URL of the file
            file_name (str): Name of the file
            user_input (str): User input to the agent
            collection_id (str): Collection ID to save the file to
            thinking_id (str): Thinking ID for activity logging
            save_to_memory (bool): Whether to save file content to agent memories for RAG.
                                   Set to True for learn endpoints, False for chat completions.

        Returns:
            str: Response from the agent
        """
        logging.info(
            f"learn_from_file called: file_url={file_url}, file_name={file_name}"
        )
        file_content = ""
        if file_name == "":
            file_name = file_url.split("/")[-1]
        if file_url.startswith(self.outputs):
            folder_path = file_url.split(f"{self.outputs}/")[1]
            file_path = os.path.normpath(
                os.path.join(self.agent_workspace, folder_path)
            )
            abs_workspace = os.path.abspath(self.agent_workspace)
            abs_file_path = os.path.abspath(file_path)
            # Ensure file path stays within the workspace directory
            if not abs_file_path.startswith(abs_workspace + os.sep):
                raise Exception(
                    "Invalid file path: attempt to access outside of the workspace."
                )
            file_path = abs_file_path
        else:
            file_data = await self.download_file_to_workspace(
                url=file_url, file_name=file_name
            )
            self.conversation.increment_attachment_count()
            if file_data == {}:
                self.conversation.log_interaction(
                    role=self.agent_name,
                    message=f"[ACTIVITY][ERROR] I was unable to read the file from {file_url}.",
                )
                return f"Unable to read the file from {file_url} ."
            file_name = file_data["file_name"]
            file_path = os.path.normpath(
                os.path.join(self.agent_workspace, collection_id, file_name)
            )
            abs_workspace = os.path.abspath(self.agent_workspace)
            abs_file_path = os.path.abspath(file_path)
            # Ensure file path stays within the workspace directory
            if not abs_file_path.startswith(abs_workspace + os.sep):
                raise Exception(
                    "Invalid file path: attempt to access outside of the workspace."
                )
            file_path = abs_file_path
        file_type = file_name.split(".")[-1]
        action_verb = "Learning" if save_to_memory else "Saving"
        action_location = "to memory" if save_to_memory else "to workspace"
        if file_type not in ["jpg", "jpeg", "png", "gif"]:
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[SUBACTIVITY][{thinking_id}] {action_verb} [{file_name}]({file_url}) {action_location}.",
            )
        if file_type in ["ppt", "pptx"]:
            # Extract text directly from PowerPoint using python-pptx
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[SUBACTIVITY][{thinking_id}] {action_verb} PowerPoint file [{file_name}]({file_url}) {action_location}.",
            )
            try:
                prs = Presentation(file_path)
                pptx_content = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    if slide_text:
                        pptx_content.append(
                            f"Slide {slide_num}:\n" + "\n".join(slide_text)
                        )
                content = "\n\n".join(pptx_content)
                pptx_content_str = (
                    f"Content from PowerPoint uploaded named `{file_name}`:\n{content}"
                )
                file_content += pptx_content_str
                self.input_tokens += get_tokens(content)
                if save_to_memory:
                    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    await self.file_reader.write_text_to_memory(
                        user_input=user_input,
                        text=f"Content from PowerPoint uploaded at {timestamp} named `{file_name}`:\n{content}",
                        external_source=f"file {file_path}",
                    )
                response = f"{'Learned' if save_to_memory else 'Saved'} [{file_name}]({file_url}) {'to memory' if save_to_memory else 'to workspace'}."
            except Exception as e:
                logging.error(f"Error reading PowerPoint file: {e}")
                return f"Failed to read PowerPoint file [{file_name}]({file_url}). Error: {str(e)}"
        if user_input == "":
            user_input = "Describe each stage of this image."
        disallowed_types = ["exe", "bin", "rar"]
        if file_type in ["ppt", "pptx"]:
            # Already handled above, response is set
            pass
        elif file_type in disallowed_types:
            response = f"[ERROR] I was unable to read the file called `{file_name}`."
        elif file_type == "pdf":
            file_content += f"Content from PDF uploaded named `{file_name}`:\n"
            with pdfplumber.open(file_path) as pdf:
                content = "\n".join([page.extract_text() for page in pdf.pages])
                file_content += content
            if "pdf_vision" in self.agent_settings:
                if (
                    self.agent_settings["pdf_vision"] != "None"
                    and self.agent_settings["pdf_vision"] != ""
                    and self.agent_settings["pdf_vision"] != None
                    and str(self.agent_settings["pdf_vision"]).lower() != "false"
                    and self.agent_settings["vision_provider"] != "None"
                    and self.agent_settings["vision_provider"] != ""
                    and self.agent_settings["vision_provider"] != None
                ):
                    with pdfplumber.open(file_path) as pdf:
                        for i, page in enumerate(pdf.pages):
                            page_image = page.to_image(width=256, height=256)
                            image_path = (
                                f"{file_path.replace('.pdf', f'_page_{i}.png')}"
                            )
                            page_image.save(image_path)
                            # Read image and encode as base64 for vision inference
                            with open(image_path, "rb") as img_file:
                                image_data = img_file.read()
                                base64_image = base64.b64encode(image_data).decode(
                                    "utf-8"
                                )
                            base64_image = f"data:image/png;base64,{base64_image}"
                            vision_response = await self.agent.vision_inference(
                                prompt=content, images=[base64_image]
                            )
                        file_content += f"Visual description from viewing uploaded PDF called `{file_name}` from page {i} with OCR:\n"
                        file_content += vision_response
            self.input_tokens += get_tokens(content)
            if save_to_memory:
                timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                await self.file_reader.write_text_to_memory(
                    user_input=user_input,
                    text=f"Content from PDF uploaded at {timestamp} named `{file_name}`:\n{content}",
                    external_source=f"file {file_path}",
                )
            response = f"{'Learned' if save_to_memory else 'Saved'} [{file_name}]({file_url}) {'to memory' if save_to_memory else 'to workspace'}."
        elif file_path.endswith(".zip"):
            extracted_zip_folder_name = f"extracted_{file_name.replace('.zip', '_zip')}"
            new_folder = os.path.normpath(
                os.path.join(
                    self.agent_workspace, collection_id, extracted_zip_folder_name
                )
            )
            file_content += f"Content from the zip file uploaded named `{file_name}`:\n"
            if new_folder.startswith(self.agent_workspace):
                with zipfile.ZipFile(file_path, "r") as zipObj:
                    zipObj.extractall(path=new_folder)
                # Build folder structure summary instead of processing each file
                # This prevents spamming activities for large repos like Flipper-IRDB
                folder_structure = []
                file_count = 0
                dir_count = 0
                for root, dirs, files in os.walk(new_folder):
                    rel_path = os.path.relpath(root, new_folder)
                    if rel_path == ".":
                        rel_path = ""
                    dir_count += len(dirs)
                    file_count += len(files)
                    # Only show top-level structure and first 2 levels to avoid overwhelming output
                    depth = rel_path.count(os.sep) if rel_path else 0
                    if depth <= 2:
                        if rel_path:
                            folder_structure.append(f"📁 {rel_path}/")
                        # Show up to 10 files per directory at depth <= 1
                        if depth <= 1 and files:
                            shown_files = files[:10]
                            for f in shown_files:
                                folder_structure.append(f"   📄 {f}")
                            if len(files) > 10:
                                folder_structure.append(
                                    f"   ... and {len(files) - 10} more files"
                                )

                structure_summary = "\n".join(
                    folder_structure[:100]
                )  # Limit to 100 lines
                if len(folder_structure) > 100:
                    structure_summary += f"\n... and more ({dir_count} total directories, {file_count} total files)"
                else:
                    structure_summary += (
                        f"\n\nTotal: {dir_count} directories, {file_count} files"
                    )

                file_content += f"\n\n**EXTRACTED ZIP CONTENTS** (use the extracted folder, NOT the .zip file):\n"
                file_content += (
                    f"Extracted folder path: `{extracted_zip_folder_name}/`\n\n"
                )
                file_content += f"Folder structure:\n{structure_summary}\n"
                file_content += f"\n**IMPORTANT**: To read files, use paths like `{extracted_zip_folder_name}/subfolder/file.ext` - do NOT try to read the .zip file directly."
                response = f"Extracted zip file [{file_name}]({file_url}) to `{extracted_zip_folder_name}/` ({file_count} files in {dir_count} directories). Use the EXTRACTED FOLDER to browse files, not the zip."
            else:
                response = (
                    f"[ERROR] I was unable to read the file called `{file_name}`."
                )
        elif file_path.endswith(".doc") or file_path.endswith(".docx"):
            content = docx2txt.process(file_path)
            docx_content = (
                f"Content from the document uploaded named `{file_name}`:\n{content}"
            )
            file_content += docx_content
            self.input_tokens += get_tokens(content)
            if save_to_memory:
                await self.file_reader.write_text_to_memory(
                    user_input=user_input,
                    text=docx_content,
                    external_source=f"file {file_path}",
                )
            response = f"{'Learned' if save_to_memory else 'Saved'} [{file_name}]({file_url}) {'to memory' if save_to_memory else 'to workspace'}."
        elif file_type == "xlsx" or file_type == "xls" or file_type == "csv":
            response, content = await self.learn_spreadsheet(
                user_input=user_input,
                file_path=file_path,
                thinking_id=thinking_id,
                save_to_memory=save_to_memory,
            )
            file_content += content
            if save_to_memory:
                await self.file_reader.write_text_to_memory(
                    user_input=user_input,
                    text=content,
                    external_source=f"file {file_path}",
                )
        elif (
            file_type == "wav"
            or file_type == "mp3"
            or file_type == "ogg"
            or file_type == "m4a"
            or file_type == "flac"
            or file_type == "wma"
            or file_type == "aac"
        ):
            audio = AudioSegment.from_file(file_path)
            audio.export(file_path, format="wav")
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[ACTIVITY] Transcribing audio file [{file_name}]({file_url}).",
            )
            audio_response = await self.audio_to_text(audio_path=file_path)
            file_content += (
                f"Transcription from the audio file uploaded named `{file_name}`:\n"
            )
            file_content += audio_response
            self.input_tokens += get_tokens(audio_response)
            if save_to_memory:
                await self.file_reader.write_text_to_memory(
                    user_input=user_input,
                    text=f"Transcription from the audio file called `{file_name}`:\n{audio_response}\n",
                    external_source=f"audio {file_name}",
                )
            response = f"Transcribed audio from [{file_name}]({file_url}) and {'saved to memory' if save_to_memory else 'saved to workspace'}."
        # If it is an image, generate a description
        elif file_type in [
            "jpg",
            "jpeg",
            "png",
            "gif",
            "webp",
            "tiff",
            "bmp",
            "svg",
        ]:
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[SUBACTIVITY] Uploaded `{file_name}` ![Uploaded {file_name}]({file_url})",
            )
            # Separate vision inference from memory writing to avoid false errors
            vision_response = None
            vision_prompt = f"The assistant has an image in context\nThe user's last message was: {user_input}\nThe uploaded image is `{file_name}`.\n\nAnswer anything relevant to the image that the user is questioning if anything, additionally, describe the image in detail."
            self.input_tokens += get_tokens(vision_prompt)
            # Read image and encode as base64 for vision inference
            with open(file_path, "rb") as img_file:
                image_data = img_file.read()
                base64_image = base64.b64encode(image_data).decode("utf-8")
            base64_image = f"data:image/{file_type};base64,{base64_image}"
            vision_response = await self.agent.vision_inference(
                prompt=vision_prompt, images=[base64_image]
            )
            # Check if vision_inference returned an error response
            if vision_response and "Unable to process request" in vision_response:
                logging.error(f"Vision inference returned error for {file_name}")
                vision_response = None

            if vision_response:
                file_content += f"Visual description from viewing uploaded image called `{file_name}`:\n"
                file_content += vision_response
                response = f"{'Learned' if save_to_memory else 'Processed'} [{file_name}]({file_url}) {'to memory' if save_to_memory else 'to workspace'}."
            else:
                response = f"I was unable to view the image called `{file_name}`. I will need to try the `View Image` ability."
        else:
            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            abs_workspace = os.path.abspath(self.agent_workspace)
            fp = os.path.abspath(os.path.normpath(file_path))
            if fp.startswith(abs_workspace + os.sep):
                try:
                    with open(fp, "r") as f:
                        content = f.read()
                except:
                    with open(fp, "rb") as f:
                        content = f.read()
                    content = base64.b64encode(content).decode("utf-8")
                file_content += (
                    f"Content from file uploaded named `{file_name}` at {timestamp}:\n"
                )
                file_content += content
                self.input_tokens += get_tokens(content)
                # Check how many lines are in the file content
                if save_to_memory:
                    lines = content.split("\n")
                    if len(lines) > 1:
                        for line_number, line in enumerate(lines):
                            await self.file_reader.write_text_to_memory(
                                user_input=user_input,
                                text=f"Content from file uploaded named `{file_name}` at {timestamp} on line number {line_number + 1}:\n{line}",
                                external_source=f"file {fp}",
                            )
                    else:
                        await self.file_reader.write_text_to_memory(
                            user_input=user_input,
                            text=f"Content from file uploaded named `{file_name}` at {timestamp}:\n{content}",
                            external_source=f"file {fp}",
                        )
                response = f"{'Learned' if save_to_memory else 'Saved'} [{file_name}]({file_url}) {'to memory' if save_to_memory else 'to workspace'}."
            else:
                response = (
                    f"[ERROR] I was unable to read the file called `{file_name}`."
                )
        self.conversation.log_interaction(
            role=self.agent_name,
            message=(
                f"[SUBACTIVITY][{thinking_id}] {response}"
                if "[ERROR]" not in response
                else f"[ACTIVITY]{response}"
            ),
        )

        # For learn endpoints (save_to_memory=True), return the simple response message
        # For chat completions (save_to_memory=False), return metadata and instructions
        if save_to_memory:
            # Learning to memory - return simple success message
            return response

        # For images with vision content, return the vision description (it's already concise)
        # The vision response is essential context for the agent to respond about the image
        image_types = ["jpg", "jpeg", "png", "gif", "webp", "tiff", "bmp", "svg"]
        if file_type in image_types and file_content:
            logging.info(
                f"Image {file_name} - returning vision response content ({get_tokens(file_content)} tokens)"
            )
            return file_content

        # Chat completions - return metadata and instructions instead of full file content
        # This keeps context manageable and the agent can use commands to access files
        if file_content:
            file_tokens = get_tokens(file_content)
            logging.info(f"learn_from_file: file={file_name}, tokens={file_tokens}")

            # Get file size
            try:
                file_size = os.path.getsize(file_path)
            except:
                file_size = len(file_content.encode("utf-8"))

            # Determine the converted file name for spreadsheets
            converted_file = None
            if file_type in ["xlsx", "xls"]:
                base_name = file_name.rsplit(".", 1)[0]
                try:
                    # Securely determine the directory for file_path
                    dir_path = os.path.normpath(os.path.dirname(file_path))
                    if not dir_path.startswith(self.agent_workspace):
                        raise Exception(
                            "Access to directory outside agent workspace is not allowed."
                        )
                    csv_files = [
                        f
                        for f in os.listdir(dir_path)
                        if f.startswith(base_name) and f.endswith(".csv")
                    ]
                    if csv_files:
                        converted_file = (
                            csv_files[0]
                            if len(csv_files) == 1
                            else f"{base_name}_*.csv (multiple sheets)"
                        )
                except Exception as e:
                    log_silenced_exception(
                        e, f"_process_file: listing CSV files for {file_name}"
                    )

            # Return only metadata and instructions - not the actual content
            instructions = self._get_file_access_instructions(
                file_name=file_name,
                file_url=file_url,
                file_path=file_path,
                file_size_bytes=file_size,
                file_tokens=file_tokens,
                converted_file_name=converted_file,
                file_type=file_type,
            )

            logging.info(
                f"File {file_name} - returning metadata only ({file_tokens} tokens content not included in context)"
            )
            return instructions

        return file_content

    async def _process_file_type_message(self, msg: dict, files: list) -> None:
        """
        Process a message content item with type "file" and nested file object.

        Handles the OpenAI-style file upload format:
        {"type": "file", "file": {"filename": "...", "file_data": "data:..."}}

        Args:
            msg: Message content item dictionary
            files: List to append downloaded file info to
        """
        if not (isinstance(msg, dict) and msg.get("type") == "file" and "file" in msg):
            return

        file_info = msg["file"]
        file_name = file_info.get("filename", f"{uuid.uuid4().hex}.txt")
        file_data_str = file_info.get("file_data", "")

        if file_data_str:
            # Download the file to workspace and add to files list
            downloaded_file = await self.download_file_to_workspace(
                url=file_data_str,
                file_name=file_name,
            )
            if downloaded_file != {}:
                files.append(downloaded_file)
                logging.info(f"Processed file upload: {file_name}")

    async def download_file_to_workspace(
        self, url: str, file_name: str = "", download_headers={}
    ):
        """
        Download a file from a URL to the workspace

        Args:
            url (str): URL of the file
            file_name (str): Name of the file

        Returns:
            str: URL of the downloaded file
        """
        try:
            # Check if URL already points to the same conversation workspace - skip re-downloading
            if self.outputs and self.conversation_id:
                workspace_prefix = f"{self.outputs}/{self.conversation_id}/"
                if url.startswith(workspace_prefix):
                    # File already exists in this conversation's workspace, no need to re-download
                    existing_file_name = url[len(workspace_prefix) :]
                    logging.info(
                        f"File already in workspace, skipping download: {existing_file_name}"
                    )
                    return {"file_name": existing_file_name, "file_url": url}

            if url.startswith("data:"):
                file_type = url.split(",")[0].split("/")[1].split(";")[0]
            else:
                if "?" in url:
                    file_type = url.split("?")[0]
                    file_type = file_type.split(".")[-1]
                else:
                    file_type = url.split(".")[-1]
            if not file_type:
                file_type = "txt"
            self.conversation_id
            if not self.conversation_id:
                self.conversation_id = self.conversation.get_conversation_id()
            file_name = (
                f"{uuid.uuid4().hex}.{file_type}" if file_name == "" else file_name
            )
            file_name = "".join(c if c.isalnum() else "_" for c in file_name)
            file_extension = file_name.split("_")[-1]
            file_name = file_name.replace(f"_{file_extension}", f".{file_extension}")
            full_path = os.path.normpath(
                os.path.join(self.conversation_workspace, file_name)
            )
            if not full_path.startswith(self.conversation_workspace):
                raise Exception("Path given not allowed")
            if "," in url:
                file_type = url.split(",")[0].split("/")[1].split(";")[0]
                file_data = base64.b64decode(url.split(",")[1])
            else:
                file_type = file_name.split(".")[-1]
                # Download the file
                try:
                    if not url.startswith("http"):
                        return {}
                    if url in ["", None]:
                        return {}
                    # SSRF protection: validate URL before making request
                    if not is_safe_url(url):
                        logging.error(f"SSRF protection blocked download from: {url}")
                        return {}
                    file_download = requests.get(url, timeout=30)
                    file_data = file_download.content
                except Exception as e:
                    logging.error(f"Error downloading file: {e}")
                    return {}
            full_path = os.path.normpath(
                os.path.join(self.conversation_workspace, file_name)
            )
            if not full_path.startswith(self.conversation_workspace):
                raise Exception("Path given not allowed")
            with open(full_path, "wb") as f:
                f.write(file_data)
            url = f"{self.outputs}/{self.conversation_id}/{file_name}"
            return {"file_name": file_name, "file_url": url}
        except Exception as e:
            logging.error(f"Error in download_file_to_workspace: {e}")
            return {}

    async def plan_task(
        self,
        user_input: str,
        websearch: bool = False,
        websearch_depth: int = 3,
        log_user_input: bool = True,
        log_output: bool = True,
        enable_new_command: bool = True,
    ):
        """
        Plan a task from a user input, create and enable a new command to execute the plan

        Args:
        user_input (str): User input to the agent
        websearch (bool): Whether to include web research in the chain
        websearch_depth (int): Depth of web research to include
        log_user_input (bool): Whether to log the user input
        log_output (bool): Whether to log the output
        enable_new_command (bool): Whether to enable the new command for the agent

        Returns:
        str: The name of the created chain
        """
        if log_user_input:
            self.conversation.log_interaction(
                role="USER",
                message=user_input,
            )
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[ACTIVITY] Determining primary objective.",
        )
        # Who is the expert here?
        # Run the prompt "Expert Determination" with the user input
        expert = await self.inference(
            user_input=user_input,
            prompt_category="Default",
            prompt_name="Expert Determination",
            log_output=False,
            log_user_input=False,
        )
        # Use the prompt "Prompt Generator" to generate a prompt for the user to provide the primary objective
        primary_objective = await self.inference(
            user_input=user_input,
            job_title=expert,
            task=user_input,
            prompt_name="Prompt Generator",
            log_output=False,
            log_user_input=False,
        )

        # primary_objective = Step 1, execute chain "Smart Prompt" with the user input to get Primary Objective
        chain_name = await self.inference(
            user_input=user_input,
            introduction=primary_objective,
            prompt_category="Default",
            prompt_name="Title a Chain",
            log_output=False,
            log_user_input=False,
        )
        chain_title = await self.convert_to_model(
            input_string=chain_name,
            model=ChainCommandName,
        )
        chain_name = chain_title.command_name
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[ACTIVITY] Breaking objective into a list of tasks.",
        )
        # numbered_list_of_tasks = Step 2, Execute prompt "Break into Steps" with `introduction` being step 1 response, websearch true if researching
        # Note - Should do this more than once to get a better list of tasks
        numbered_list_of_tasks = await self.inference(
            user_input=user_input,
            introduction=primary_objective,
            prompt_category="Default",
            prompt_name="Break into Steps",
            websearch=websearch,
            websearch_depth=websearch_depth,
            injected_memories=100,
            log_output=False,
            log_user_input=False,
        )
        task_list = await self.convert_to_model(
            input_string=numbered_list_of_tasks,
            model=TasksToDo,
        )
        self.chain.add_chain(chain_name=chain_name)
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[ACTIVITY] Creating new command `{chain_name}`.",
        )
        i = 1
        total_tasks = len(task_list.tasks)
        x = 1
        # First step in the chain should be to disable the command so that the agent doesn't try to execute it while executing it
        self.chain.add_chain_step(
            chain_name=chain_name,
            agent_name=self.agent_name,
            step_number=i,
            prompt_type="Command",
            prompt={
                "command_name": "Disable Command",
                "command_args": {
                    "agent_name": self.agent_name,
                    "command_name": chain_name,
                },
            },
        )
        i += 1
        for task in task_list.tasks:
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[ACTIVITY] Planning task `{x}` of `{total_tasks}`.",
            )
            x += 1
            # Create a smart prompt with the objective and task in context
            self.chain.add_chain_step(
                chain_name=chain_name,
                agent_name=self.agent_name,
                step_number=i,
                prompt_type="Chain",
                prompt={
                    "chain_name": "Smart Prompt",
                    "input": f"Primary Objective to keep in mind while working on the task: {primary_objective} \nThe only task to complete to move towards the objective: {task}",
                },
            )
            i += 1
            self.chain.add_chain_step(
                chain_name=chain_name,
                agent_name=self.agent_name,
                step_number=i,
                prompt_type="Chain",
                prompt={
                    "chain": (
                        "Smart Instruct"
                        if websearch
                        else "Smart Instruct - No Research"
                    ),
                    "input": "{STEP" + str(i - 1) + "}",
                },
            )
            i += 1
        list_of_tasks = "\n".join(
            [f"{i}. {task}" for i, task in enumerate(task_list.tasks, 1)]
        )
        # Enable the command of the chain name
        if enable_new_command:
            self.agent.update_agent_config(
                new_config={chain_name: True}, config_key="commands"
            )
            message = f"I have created a new command called `{chain_name}`. The tasks will be executed in the following order:\n{list_of_tasks}\n\nWould you like me to execute `{chain_name}` now?"
        else:
            message = f"I have created a new command called `{chain_name}`. The tasks will be executed in the following order:\n{list_of_tasks}\n\nIf you are able to enable the command, I can execute it for you. Alternatively, you can execute the command manually."
        if log_output:
            self.conversation.log_interaction(
                role=self.agent_name,
                message=message,
            )
        return {
            "chain_name": chain_name,
            "message": message,
            "tasks": list_of_tasks,
        }

    async def update_planned_task(
        self,
        chain_name: str,
        user_input: str,
        log_user_input: bool = True,
        log_output: bool = True,
        enable_new_command: bool = True,
    ):
        """
        Modify the chain based on user input

        Args:
        chain_name (str): Name of the chain to update
        user_input (str): User input to the agent
        log_user_input (bool): Whether to log the user input
        log_output (bool): Whether to log the output

        Returns:
        str: Response from the agent
        """
        # Basically just delete the old chain after we extract the tasks and then run the plan_task function with more input from the user.
        current_chain = self.chain.get_chain(chain_name=chain_name)
        # This function is still a work in progress
        # Need to

        self.chain.delete_chain(chain_name=chain_name)
        return await self.plan_task(
            user_input=user_input,
            log_user_input=log_user_input,
            log_output=log_output,
            enable_new_command=enable_new_command,
        )

    def remove_tagged_content(self, text: str, tag: str) -> str:
        """Safely remove content between tags without using regex."""
        start_tag = f"[{tag}]"
        end_tag = f"[/{tag}]"

        while True:
            start = text.find(start_tag)
            if start == -1:
                break
            end = text.find(end_tag, start)
            if end == -1:
                break
            text = text[:start] + text[end + len(end_tag) :]
        return text

    async def chat_completions(self, prompt: ChatCompletions):
        """
        Generate an OpenAI style chat completion response with a ChatCompletion prompt

        Args:
            prompt (ChatCompletions): Chat completions prompt

        Returns:
            dict: Chat completion response
        """
        # conversation_name = prompt.user
        c = self.conversation
        conversation_id = self.conversation_id

        # Register this conversation as active
        task = asyncio.current_task()
        worker_registry.register_conversation(
            conversation_id=conversation_id,
            user_id=self.auth.user_id,
            agent_name=self.agent_name,
            task=task,
        )

        try:
            return await self._execute_chat_completions(prompt)
        except asyncio.CancelledError:
            logging.info(
                f"Chat completion cancelled for conversation {conversation_id}"
            )
            raise
        except Exception as e:
            logging.error(
                f"Error in chat completions for conversation {conversation_id}: {e}"
            )
            raise
        finally:
            # Always unregister when done
            worker_registry.unregister_conversation(conversation_id)

    async def _execute_chat_completions(self, prompt: ChatCompletions):
        """
        Internal method that does the actual chat completion processing
        """
        # Validate that messages is provided and not empty
        if not prompt.messages:
            raise ValueError(
                "The 'messages' field is required and must contain at least one message."
            )
        c = self.conversation
        conversation_id = self.conversation_id
        urls = []
        files = []
        new_prompt = ""
        browse_links = True
        tts = False
        websearch = False
        language = "en"
        log_output = True
        log_user_input = True
        command_overrides = None
        if prompt.tools:
            command_overrides = prompt.tools
        if "websearch" in self.agent_settings:
            websearch = str(self.agent_settings["websearch"]).lower() == "true"
        if "mode" in self.agent_settings:
            mode = self.agent_settings["mode"]
        else:
            mode = "prompt"
        if "prompt_name" in self.agent_settings:
            prompt_name = self.agent_settings["prompt_name"]
        else:
            prompt_name = "Think About It"
        if "prompt_category" in self.agent_settings:
            prompt_category = self.agent_settings["prompt_category"]
        else:
            prompt_category = "Default"
        if "LANGUAGE" in self.agent_settings:
            language = str(self.agent_settings["LANGUAGE"]).lower()
        prompt_args = {}
        if "prompt_args" in self.agent_settings:
            prompt_args = (
                json.loads(self.agent_settings["prompt_args"])
                if isinstance(self.agent_settings["prompt_args"], str)
                else self.agent_settings["prompt_args"]
            )
        if "context_results" in self.agent_settings:
            context_results = int(self.agent_settings["context_results"])
        else:
            context_results = 5
        if "injected_memories" in self.agent_settings:
            context_results = int(self.agent_settings["injected_memories"])
        if "conversation_results" in self.agent_settings:
            conversation_results = int(self.agent_settings["conversation_results"])
        else:
            conversation_results = 6
        if "command_name" in self.agent_settings:
            command_name = self.agent_settings["command_name"]
        else:
            command_name = ""
        if "command_args" in self.agent_settings:
            try:
                command_args = (
                    json.loads(self.agent_settings["command_args"])
                    if isinstance(self.agent_settings["command_args"], str)
                    else self.agent_settings["command_args"]
                )
            except Exception as e:
                command_args = {}
        else:
            command_args = {}
        if "command_variable" in self.agent_settings:
            command_variable = self.agent_settings["command_variable"]
        else:
            command_variable = "text"
        if "chain_name" in self.agent_settings:
            chain_name = self.agent_settings["chain_name"]
        else:
            chain_name = ""
        if "chain_args" in self.agent_settings:
            chain_args = (
                json.loads(self.agent_settings["chain_args"])
                if isinstance(self.agent_settings["chain_args"], str)
                else self.agent_settings["chain_args"]
            )
        else:
            chain_args = {}
        if "tts_provider" in self.agent_settings:
            tts_provider = str(self.agent_settings["tts_provider"]).lower()
            if tts_provider != "none" and tts_provider != "":
                if "tts" in self.agent_settings:
                    tts = str(self.agent_settings["tts"]).lower() == "true"
        analyze_user_input = False
        if "analyze_user_input" in self.agent_settings:
            analyze_user_input = (
                str(self.agent_settings["analyze_user_input"]).lower() == "true"
            )
        include_sources = False
        if "include_sources" in self.agent_settings:
            include_sources = (
                str(self.agent_settings["include_sources"]).lower() == "true"
            )
        disable_commands = False
        running_command = None
        additional_context = ""
        parent_activity_id = None
        has_tool_result = False  # Track if this is a tool result continuation
        tool_result_text = ""  # Store tool result text separately
        for message in prompt.messages:
            if "mode" in message:
                if message["mode"] in ["prompt", "command", "chain"]:
                    mode = message["mode"]
            if "log_output" in message:
                log_output = str(message["log_output"]).lower() == "true"
            if "log_user_input" in message:
                log_user_input = str(message["log_user_input"]).lower() == "true"
            if "injected_memories" in message:
                context_results = int(message["injected_memories"])
            if "parent_activity_id" in message:
                parent_activity_id = message["parent_activity_id"]
            if "language" in message:
                language = message["language"]
            if "conversation_results" in message:
                conversation_results = int(message["conversation_results"])
            if "prompt_category" in message:
                prompt_category = message["prompt_category"]
            if "prompt_name" in message:
                prompt_name = message["prompt_name"]
            if "prompt_args" in message:
                prompt_args = (
                    json.loads(message["prompt_args"])
                    if isinstance(message["prompt_args"], str)
                    else message["prompt_args"]
                )
            if "command_name" in message:
                command_name = message["command_name"]
            if "command_args" in message:
                command_args = (
                    json.loads(message["command_args"])
                    if isinstance(message["command_args"], str)
                    else message["command_args"]
                )
            if "command_variable" in message:
                command_variable = message["command_variable"]
            if "chain_name" in message:
                chain_name = message["chain_name"]
            if "chain_args" in message:
                chain_args = (
                    json.loads(message["chain_args"])
                    if isinstance(message["chain_args"], str)
                    else message["chain_args"]
                )
            if "browse_links" in message:
                browse_links = str(message["browse_links"]).lower() == "true"
            if "tts" in message:
                tts = str(message["tts"]).lower() == "true"
            if "websearch" in message:
                websearch = str(message["websearch"]).lower() == "true"
            if "analyze_user_input" in message:
                analyze_user_input = (
                    str(message["analyze_user_input"]).lower() == "true"
                )
            if "context" in message:
                additional_context += "\n" + str(message["context"]).strip()
            if "include_sources" in message:
                include_sources = str(message["include_sources"]).lower() == "true"
            download_headers = {}
            if "download_headers" in message:
                download_headers = (
                    json.loads(message["download_headers"])
                    if isinstance(message["download_headers"], str)
                    else message["download_headers"]
                )
            if "disable_commands" in message:
                disable_commands = str(message["disable_commands"]).lower() == "true"
            if "running_command" in message:
                running_command = message["running_command"]
            if "content" not in message:
                continue
            role = message["role"] if "role" in message else "User"
            # Handle tool result messages - treat as context for continuing conversation
            if role.lower() == "tool":
                has_tool_result = True  # Mark that this conversation has tool results
                tool_call_id = message.get("tool_call_id", "unknown")
                tool_content = message["content"]
                if isinstance(tool_content, str):
                    # Add tool result as prompt content - agent should respond based on this
                    new_prompt += f"{tool_content}\n\n"
                    tool_result_text = tool_content  # Store for logging
                    # Extract URLs from tool result text (may be JSON or plain text)
                    # Look for GitHub URLs specifically to download repos to workspace
                    # Match URLs but stop at common JSON/text delimiters including backslash
                    url_pattern = r'https?://[^\s"\'<>\}\]\,\\]+'
                    found_urls = re.findall(url_pattern, tool_content)
                    for found_url in found_urls:
                        # Clean up any trailing punctuation that got matched
                        found_url = found_url.rstrip('",}]\\:')
                        # Skip if URL looks malformed or has leftover JSON
                        if "workflow" in found_url or len(found_url) > 200:
                            continue
                        if found_url.startswith("https://github.com/"):
                            do_not_pull_repo = [
                                "/pull/",
                                "/issues",
                                "/discussions",
                                "/actions/",
                                "/projects",
                                "/security",
                                "/releases",
                                "/commits",
                                "/branches",
                                "/tags",
                                "/stargazers",
                                "/watchers",
                                "/network",
                                "/settings",
                                "/compare",
                                "/archive",
                            ]
                            if any(x in found_url for x in do_not_pull_repo):
                                urls.append(found_url)
                                logging.info(
                                    f"[chat_completions] Found GitHub non-repo URL in tool result: {found_url}"
                                )
                            else:
                                # Download the GitHub repo as zip
                                logging.info(
                                    f"[chat_completions] Found GitHub repo URL in tool result: {found_url}"
                                )
                                github_user = self.agent_settings.get("GITHUB_USERNAME")
                                github_token = self.agent_settings.get("GITHUB_TOKEN")
                                github_repo = found_url.replace(
                                    "https://github.com/", ""
                                ).replace("https://www.github.com/", "")
                                # Parse out branch from /tree/branch if present
                                tool_github_branch = "main"
                                if "/tree/" in github_repo:
                                    parts = github_repo.split("/tree/")
                                    github_repo = parts[0]
                                    if len(parts) > 1:
                                        # Branch might have subpath after it
                                        branch_and_path = parts[1].split("/")
                                        tool_github_branch = branch_and_path[0]
                                user_parts = github_repo.split("/")
                                if len(user_parts) >= 2:
                                    user = user_parts[0]
                                    repo = user_parts[1]
                                    # Clean up repo name
                                    for symbol in [
                                        " ",
                                        "\n",
                                        "\t",
                                        "\r",
                                        "\\",
                                        "/",
                                        ":",
                                        "*",
                                        "?",
                                        '"',
                                        "<",
                                        ">",
                                    ]:
                                        repo = repo.replace(symbol, "")
                                        user = user.replace(symbol, "")
                                        tool_github_branch = tool_github_branch.replace(
                                            symbol, ""
                                        )
                                    repo_url = f"https://github.com/{user}/{repo}/archive/refs/heads/{tool_github_branch}.zip"
                                    logging.info(
                                        f"[chat_completions] Downloading GitHub repo: {repo_url}"
                                    )
                                    try:
                                        if github_user and github_token:
                                            response = requests.get(
                                                repo_url,
                                                auth=(github_user, github_token),
                                            )
                                        else:
                                            response = requests.get(repo_url)
                                        if response.status_code != 200:
                                            # Try master branch
                                            tool_github_branch = "master"
                                            repo_url = f"https://github.com/{user}/{repo}/archive/refs/heads/{tool_github_branch}.zip"
                                            if github_user and github_token:
                                                response = requests.get(
                                                    repo_url,
                                                    auth=(github_user, github_token),
                                                )
                                            else:
                                                response = requests.get(repo_url)
                                        if response.status_code == 200:
                                            file_name = f"{user}_{repo}_{tool_github_branch}.zip"
                                            file_path = os.path.normpath(
                                                os.path.join(
                                                    self.agent_workspace,
                                                    conversation_id,
                                                    file_name,
                                                )
                                            )
                                            # Validate path stays within workspace to prevent path traversal
                                            abs_workspace = os.path.abspath(
                                                self.agent_workspace
                                            )
                                            abs_file_path = os.path.abspath(file_path)
                                            if not abs_file_path.startswith(
                                                abs_workspace + os.sep
                                            ):
                                                logging.error(
                                                    f"[chat_completions] Invalid file path: attempt to access outside workspace"
                                                )
                                                continue
                                            os.makedirs(
                                                os.path.dirname(abs_file_path),
                                                exist_ok=True,
                                            )
                                            with open(abs_file_path, "wb") as f:
                                                f.write(response.content)
                                            logging.info(
                                                f"[chat_completions] Downloaded GitHub repo to: {file_path}"
                                            )
                                            # Append as dict with file_name and file_url for consistency
                                            file_url = f"{self.outputs}/{conversation_id}/{file_name}"
                                            files.append(
                                                {
                                                    "file_name": file_name,
                                                    "file_url": file_url,
                                                }
                                            )
                                    except Exception as e:
                                        logging.error(
                                            f"[chat_completions] Failed to download GitHub repo: {e}"
                                        )
                        elif found_url.startswith("http"):
                            # Add other URLs to the browse list
                            urls.append(found_url)
                elif isinstance(tool_content, list):
                    # Handle multipart tool results (text + images)
                    for part in tool_content:
                        if isinstance(part, dict):
                            if "text" in part:
                                new_prompt += f"{part['text']}\n\n"
                            # Process file type parts
                            await self._process_file_type_message(part, files)
                            # Process image_url type parts (from ESP32 camera, etc)
                            if part.get("type") == "image_url" and "image_url" in part:
                                img_url_data = part["image_url"]
                                url = (
                                    img_url_data.get("url", "")
                                    if isinstance(img_url_data, dict)
                                    else str(img_url_data)
                                )
                                if url:
                                    agent_id = (
                                        self.agent.agent_id if self.agent else None
                                    )
                                    # Check if this is a workspace URL missing agent_id
                                    # Pattern: /outputs/{conversation_id}/{filename} -> needs agent_id
                                    if "/outputs/" in url and agent_id:
                                        # Extract path after /outputs/
                                        outputs_idx = url.find("/outputs/")
                                        path_after_outputs = url[
                                            outputs_idx + 9 :
                                        ]  # After "/outputs/"
                                        parts = path_after_outputs.split("/")
                                        # If only 2 parts (conversation_id/filename), insert agent_id
                                        if (
                                            len(parts) == 2
                                            and parts[0] == self.conversation_id
                                        ):
                                            # Reconstruct with agent_id
                                            base_url = url[:outputs_idx]
                                            url = f"{base_url}/outputs/{agent_id}/{path_after_outputs}"
                                            logging.info(
                                                f"[chat_completions] Fixed workspace URL: {url}"
                                            )
                                    # Download to workspace
                                    file_name = (
                                        url.split("/")[-1]
                                        if "/" in url
                                        else f"{uuid.uuid4().hex}.jpg"
                                    )
                                    downloaded = await self.download_file_to_workspace(
                                        url=url, file_name=file_name
                                    )
                                    if downloaded and downloaded != {}:
                                        files.append(downloaded)
                                        logging.info(
                                            f"[chat_completions] Downloaded tool image: {file_name}"
                                        )
                logging.info(
                    f"[chat_completions] Processed tool result for {tool_call_id}"
                )
                continue
            if isinstance(message["content"], str):
                if role.lower() == "system":
                    if "/" in message["content"]:
                        new_prompt += f"{message['content']}\n\n"
                if role.lower() == "user":
                    new_prompt += f"{message['content']}\n\n"
            if isinstance(message["content"], list):
                for msg in message["content"]:
                    if "text" in msg:
                        if role.lower() == "user":
                            new_prompt += f"{msg['text']}\n\n"
                    # Process file type messages (non-streaming)
                    await self._process_file_type_message(msg, files)
                    # Iterate over the msg to find _url in one of the keys then use the value of that key unless it has a "url" under it
                    if isinstance(msg, dict):
                        for key, value in msg.items():
                            if "_url" in key:
                                url = str(value["url"] if "url" in value else value)
                                if url.startswith("https://github.com/"):
                                    do_not_pull_repo = [
                                        "/pull/",
                                        "/issues",
                                        "/discussions",
                                        "/actions/",
                                        "/projects",
                                        "/security",
                                        "/releases",
                                        "/commits",
                                        "/branches",
                                        "/tags",
                                        "/stargazers",
                                        "/watchers",
                                        "/network",
                                        "/settings",
                                        "/compare",
                                        "/archive",
                                    ]
                                    if any(x in url for x in do_not_pull_repo):
                                        # If the URL is not a repository, don't pull it
                                        urls.append(url)
                                    else:
                                        # Download the zip for the repo
                                        github_user = (
                                            self.agent_settings["GITHUB_USERNAME"]
                                            if "GITHUB_USERNAME" in self.agent_settings
                                            else None
                                        )
                                        github_token = (
                                            self.agent_settings["GITHUB_TOKEN"]
                                            if "GITHUB_TOKEN" in self.agent_settings
                                            else None
                                        )
                                        github_repo = url.replace(
                                            "https://github.com/", ""
                                        )
                                        github_repo = github_repo.replace(
                                            "https://www.github.com/", ""
                                        )
                                        if not github_branch:
                                            github_branch = "main"
                                        user = github_repo.split("/")[0]
                                        repo = github_repo.split("/")[1]
                                        if " " in repo:
                                            repo = repo.split(" ")[0]
                                        if "\n" in repo:
                                            repo = repo.split("\n")[0]
                                        # Remove any symbols that would not be in the user, repo, or branch
                                        for symbol in [
                                            " ",
                                            "\n",
                                            "\t",
                                            "\r",
                                            "\\",
                                            "/",
                                            ":",
                                            "*",
                                            "?",
                                            '"',
                                            "<",
                                            ">",
                                        ]:
                                            repo = repo.replace(symbol, "")
                                            user = user.replace(symbol, "")
                                            github_branch = github_branch.replace(
                                                symbol, ""
                                            )
                                        repo_url = f"https://github.com/{user}/{repo}/archive/refs/heads/{github_branch}.zip"
                                        try:
                                            if github_user and github_token:
                                                response = requests.get(
                                                    repo_url,
                                                    auth=(github_user, github_token),
                                                )
                                            else:
                                                response = requests.get(repo_url)
                                        except:
                                            github_branch = "master"
                                            repo_url = f"https://github.com/{user}/{repo}/archive/refs/heads/{github_branch}.zip"
                                            try:
                                                if github_user and github_token:
                                                    response = requests.get(
                                                        repo_url,
                                                        auth=(
                                                            github_user,
                                                            github_token,
                                                        ),
                                                    )
                                                else:
                                                    response = requests.get(repo_url)
                                            except:
                                                pass
                                        if response.status_code == 200:
                                            file_name = (
                                                f"{user}_{repo}_{github_branch}.zip"
                                            )
                                            file_data = response.content
                                            file_path = os.path.normpath(
                                                os.path.join(
                                                    self.agent_workspace,
                                                    conversation_id,
                                                    file_name,
                                                )
                                            )
                                            # Validate path stays within workspace to prevent path traversal
                                            abs_workspace = os.path.abspath(
                                                self.agent_workspace
                                            )
                                            abs_file_path = os.path.abspath(file_path)
                                            if abs_file_path.startswith(
                                                abs_workspace + os.sep
                                            ):
                                                with open(abs_file_path, "wb") as f:
                                                    f.write(file_data)
                                                files.append(
                                                    {
                                                        "file_name": file_name,
                                                        "file_url": f"{self.outputs}/{conversation_id}/{file_name}",
                                                    }
                                                )
                                        else:
                                            urls.append(url)
                                if "file_name" in msg:
                                    file_name = str(msg["file_name"])
                                else:
                                    file_name = ""
                                if key != "audio_url":
                                    downloaded_file = (
                                        await self.download_file_to_workspace(
                                            url=url,
                                            file_name=file_name,
                                            download_headers=download_headers,
                                        )
                                    )
                                    if downloaded_file != {}:
                                        files.append(downloaded_file)
                                    else:
                                        c.log_interaction(
                                            role=self.agent_name,
                                            message=f"[SUBACTIVITY][{thinking_id}][ERROR] I was unable to read from the URL specified.",
                                        )
                                else:
                                    # If there is an audio_url, it is the user's voice input that needs transcribed before running inference
                                    audio_file_info = (
                                        await self.download_file_to_workspace(
                                            url=url,
                                            file_name=(
                                                file_name
                                                if file_name
                                                else "recording.wav"
                                            ),
                                        )
                                    )
                                    if (
                                        not audio_file_info
                                        or "file_name" not in audio_file_info
                                    ):
                                        logging.error(
                                            f"Failed to download audio file from URL (length: {len(url) if url else 0})"
                                        )
                                        continue
                                    full_path = os.path.normpath(
                                        os.path.join(
                                            self.agent_workspace,
                                            conversation_id,
                                            audio_file_info["file_name"],
                                        )
                                    )
                                    if not full_path.startswith(self.agent_workspace):
                                        raise Exception("Path given not allowed")
                                    audio_file_path = os.path.join(
                                        self.agent_workspace,
                                        conversation_id,
                                        audio_file_info["file_name"],
                                    )
                                    if os.path.normpath(audio_file_path).startswith(
                                        self.agent_workspace
                                    ):
                                        wav_file = os.path.join(
                                            self.agent_workspace,
                                            conversation_id,
                                            f"{uuid.uuid4().hex}.wav",
                                        )
                                        AudioSegment.from_file(
                                            audio_file_path
                                        ).set_frame_rate(16000).export(
                                            wav_file, format="wav"
                                        )
                                        transcribed_audio = await self.audio_to_text(
                                            audio_path=wav_file,
                                        )
                                        if len(transcribed_audio) < 1:
                                            return {
                                                "id": self.conversation_id,
                                                "object": "chat.completion",
                                                "created": int(time.time()),
                                                "model": self.agent_name,
                                                "choices": [
                                                    {
                                                        "index": 0,
                                                        "message": {
                                                            "role": "assistant",
                                                            "content": "No input received.",
                                                        },
                                                        "finish_reason": "stop",
                                                        "logprobs": None,
                                                    }
                                                ],
                                                "usage": {
                                                    "prompt_tokens": 0,
                                                    "completion_tokens": 0,
                                                    "total_tokens": 0,
                                                },
                                            }
                                        new_prompt += transcribed_audio
        # Save the original user prompt before adding file info (for logging)
        original_user_prompt = new_prompt.strip()
        # Add file info to the prompt (for context to the agent)
        for file in files:
            new_prompt += f"\nUploaded file: `{file['file_name']}`."
        if "log_output" in prompt_args:
            log_output = str(prompt_args["log_output"]).lower() == "true"
            del prompt_args["log_output"]
        if "log_user_input" in prompt_args:
            log_user_input = str(prompt_args["log_user_input"]).lower() == "true"
            del prompt_args["log_user_input"]
        if log_user_input and not has_tool_result:
            # Log the original user input, not the modified one with file names appended
            # Don't log tool results as USER - they should be logged as TOOL subactivity
            c.log_interaction(role="USER", message=original_user_prompt)
        thinking_id = ""
        if log_output:
            thinking_id = c.get_thinking_id(agent_name=self.agent_name)
        # Log tool result as a subactivity under the thinking_id, not as a user message
        if has_tool_result and tool_result_text and thinking_id:
            c.log_interaction(
                role=self.agent_name,
                message=f"[SUBACTIVITY][{thinking_id}] Received tool result: \n```\n{tool_result_text}\n```",
            )
        file_contents = []
        current_input_tokens = get_tokens(new_prompt)
        for file in files:
            content = await self.learn_from_file(
                file_url=file["file_url"],
                file_name=file["file_name"],
                user_input=new_prompt,
                collection_id=self.conversation_id,
                thinking_id=thinking_id,
            )
            file_contents.append(content)
        if file_contents:
            file_content = "\n".join(file_contents)
            file_tokens = get_tokens(file_content)
            current_input_tokens = file_tokens + current_input_tokens
        else:
            file_content = ""
            current_input_tokens = self.input_tokens
        if "user_input" in prompt_args:
            del prompt_args["user_input"]
        if "prompt_name" in prompt_args:
            prompt_name = prompt_args["prompt_name"]
            del prompt_args["prompt_name"]
        if "prompt_category" in prompt_args:
            prompt_category = prompt_args["prompt_category"]
            del prompt_args["prompt_category"]
        if "websearch" in prompt_args:
            websearch = prompt_args["websearch"]
            del prompt_args["websearch"]
        if "browse_links" in prompt_args:
            browse_links = prompt_args["browse_links"]
            del prompt_args["browse_links"]
        if "tts" in prompt_args:
            tts = prompt_args["voice_response"]
            del prompt_args["tts"]
        if "context_results" in prompt_args:
            context_results = prompt_args["context_results"]
            del prompt_args["context_results"]
        if "conversation_results" in prompt_args:
            conversation_results = prompt_args["conversation_results"]
            del prompt_args["conversation_results"]
        if "analyze_user_input" in prompt_args:
            analyze_user_input = prompt_args["analyze_user_input"]
            del prompt_args["analyze_user_input"]
        if "voice_response" in prompt_args:
            tts = prompt_args["voice_response"]
            del prompt_args["voice_response"]
        if "injected_memories" in prompt_args:
            context_results = prompt_args["injected_memories"]
            del prompt_args["injected_memories"]
        if "shots" in prompt_args:
            del prompt_args["shots"]
        if "data_analysis" in prompt_args:
            del prompt_args["data_analysis"]
        if disable_commands:
            prompt_args["disable_commands"] = True
        if running_command:
            prompt_args["running_command"] = running_command
        await self.learn_from_websites(
            urls=urls,
            summarize_content=False,
        )
        data_analysis = ""
        if analyze_user_input:
            data_analysis = await self.analyze_data(user_input=new_prompt)
        if mode == "command" and command_name and command_variable:
            try:
                command_args = (
                    json.loads(self.agent_settings["command_args"])
                    if isinstance(self.agent_settings["command_args"], str)
                    else self.agent_settings["command_args"]
                )
            except Exception as e:
                command_args = {}
            command_args[self.agent_settings["command_variable"]] = new_prompt
            response = await self.execute_command(
                command_name=self.agent_settings["command_name"],
                command_args=command_args,
                voice_response=tts,
            )
        elif mode == "chain" and chain_name:
            chain_name = self.agent_settings["chain_name"]
            try:
                chain_args = (
                    json.loads(self.agent_settings["chain_args"])
                    if isinstance(self.agent_settings["chain_args"], str)
                    else self.agent_settings["chain_args"]
                )
            except Exception as e:
                chain_args = {}
            response = await self.execute_chain(
                chain_name=chain_name,
                user_input=new_prompt,
                agent_override=self.agent_name,
                chain_args=chain_args,
                log_user_input=False,
                voice_response=tts,
            )
        elif mode == "prompt":
            if current_input_tokens < self.agent.max_input_tokens:
                if file_content:
                    prompt_args["uploaded_file_data"] = file_content
            if len(language) > 2:
                language = language[:2]
            if "context" in prompt_args:
                additional_context += "\n" + prompt_args["context"]
                del prompt_args["context"]
            response = await self.inference(
                user_input=new_prompt,
                prompt_name=prompt_name,
                prompt_category=prompt_category,
                injected_memories=context_results,
                conversation_results=conversation_results,
                shots=prompt.n,
                websearch=websearch,
                browse_links=browse_links,
                voice_response=tts,
                log_user_input=False,
                log_output=False,
                enable_command_selection=True,  # Enable intelligent command selection for main user interactions
                data_analysis=data_analysis,
                language=language,
                include_sources=include_sources,
                context=additional_context,
                command_overrides=command_overrides,
                parent_activity_id=parent_activity_id,
                **prompt_args,
            )
            if response.startswith(f"{self.agent_name}:"):
                response = response[len(f"{self.agent_name}:") :]
            if response.startswith(f"{self.agent_name} :"):
                response = response[len(f"{self.agent_name} :") :]
            thoughts_and_reflections = ""
            if "<answer>" in response:
                if "</answer>" not in response:
                    response += "</answer>"
                try:
                    thoughts_and_reflections = response.split("<answer>")[0]
                except:
                    thoughts_and_reflections = ""
                try:
                    after_thoughts = response.split("</answer>")[1]
                    if len(after_thoughts) > 10:
                        thoughts_and_reflections += after_thoughts
                except:
                    pass
                answer = response.split("<answer>")[-1]
                answer = answer.split("</answer>")[0]
                response = answer
            if log_output:
                if thoughts_and_reflections:
                    # Before logging the response, lets get all activities matching the `thinking_id` mermaid diagram
                    enable_mermaid = False
                    if "enable_mermaid" in self.agent_settings:
                        enable_mermaid = (
                            str(self.agent_settings["enable_mermaid"]).lower() == "true"
                        )
                    if enable_mermaid:
                        activities = c.get_subactivities(thinking_id)
                        if activities:
                            activity_prompt = f"{new_prompt}\n\n{activities}\n\nReview the detailed activities list and create a mermaid diagram that describes the paths taken during the detailed activities that were performed based on the user input. This mermaid diagram should start with ```mermaid\nContent of the diagram\n```\ninside of the <answer> block as the final response. The activities describe the thoughts in steps that ultimately led to the response from the assistant to the user based on the user input. Be as detailed as possible with the diagram. Ensure each item in the diagram is in quotes."
                            mermaid_diagram = await self.inference(
                                user_input=activity_prompt,
                                prompt_category="Default",
                                prompt_name="Think About It",
                                log_output=False,
                                log_user_input=False,
                                voice_response=False,
                                analyze_user_input=False,
                                browse_links=False,
                                websearch=False,
                                disable_commands=True,
                                conversation_name=self.conversation_name,
                            )
                            if mermaid_diagram:
                                mermaid_diagram = mermaid_diagram.split("<answer>")[
                                    -1
                                ].split("</answer>")[0]
                                c.log_interaction(
                                    role=self.agent_name,
                                    message=f"[SUBACTIVITY][{thinking_id}][DIAGRAM] Generated diagram describing thoughts.\n{mermaid_diagram}",
                                )
                    c.update_message_by_id(
                        message_id=thinking_id,
                        new_message=f"[ACTIVITY] Completed activities.",
                    )
                self.conversation.log_interaction(
                    role=self.agent_name,
                    message=response,
                )
                # Rename new conversations after response is complete
                # Run as background task so it doesn't block the response being returned
                if self.conversation_name == "-":
                    asyncio.create_task(self.rename_new_conversation(new_prompt))
        if isinstance(response, dict):
            response = json.dumps(response, indent=2)
        if not isinstance(response, str):
            response = str(response)
        try:
            prompt_tokens = get_tokens(new_prompt) + self.input_tokens
            completion_tokens = get_tokens(response)
            total_tokens = int(prompt_tokens) + int(completion_tokens)
        except:
            if not response:
                response = "Unable to retrieve response."
                logging.error(f"Error getting response: {response}")
        response = self.remove_tagged_content(response, "execute")
        response = self.remove_tagged_content(response, "output")

        # Check if there are pending remote commands (client-defined tools)
        pending_commands = getattr(
            self.agent_interactions, "_pending_remote_commands", []
        )
        if pending_commands:
            # Build tool_calls response for client-defined tools
            tool_calls = []
            for idx, cmd in enumerate(pending_commands):
                tool_call = {
                    "id": cmd.get("request_id", f"call_{idx}"),
                    "type": "function",
                    "function": {
                        "name": cmd.get("tool_name", "unknown"),
                        "arguments": json.dumps(cmd.get("tool_args", {})),
                    },
                }
                tool_calls.append(tool_call)

            logging.info(
                f"[chat_completions] Returning {len(tool_calls)} tool_calls for client-defined tools"
            )

            res_model = {
                "id": self.conversation_id,
                "object": "chat.completion",
                "created": int(time.time()),
                "model": self.agent_name,
                "choices": [
                    {
                        "index": 0,
                        "message": {
                            "role": "assistant",
                            "content": None,
                            "tool_calls": tool_calls,
                        },
                        "finish_reason": "tool_calls",
                        "logprobs": None,
                    }
                ],
                "usage": {
                    "prompt_tokens": prompt_tokens,
                    "completion_tokens": completion_tokens,
                    "total_tokens": total_tokens,
                },
            }
            # Clear the pending commands
            self.agent_interactions._pending_remote_commands = []
            return res_model

        res_model = {
            "id": self.conversation_id,
            "object": "chat.completion",
            "created": int(time.time()),
            "model": self.agent_name,
            "choices": [
                {
                    "index": 0,
                    "message": {
                        "role": "assistant",
                        "content": str(response),
                    },
                    "finish_reason": "stop",
                    "logprobs": None,
                }
            ],
            "usage": {
                "prompt_tokens": prompt_tokens,
                "completion_tokens": completion_tokens,
                "total_tokens": total_tokens,
            },
        }
        return res_model

    async def chat_completions_stream(self, prompt: ChatCompletions):
        """
        Generate a streaming OpenAI style chat completion response with a ChatCompletion prompt

        Args:
            prompt (ChatCompletions): Chat completions prompt

        Yields:
            str: Server-Sent Events formatted streaming response chunks
        """
        import json
        import time
        import asyncio

        # Validate that messages is provided and not empty
        if not prompt.messages:
            raise ValueError(
                "The 'messages' field is required and must contain at least one message."
            )

        conversation_id = self.conversation_id
        chunk_id = conversation_id  # Use conversation_id as the chunk ID
        created_time = int(time.time())

        # Register this conversation as active
        task = asyncio.current_task()
        worker_registry.register_conversation(
            conversation_id=conversation_id,
            user_id=self.auth.user_id,
            agent_name=self.agent_name,
            task=task,
        )

        try:
            # Execute the streaming chat completion
            async for chunk in self._execute_chat_completions_stream(prompt):
                yield chunk
        except asyncio.CancelledError:
            # Handle graceful stop
            final_chunk = {
                "id": chunk_id,
                "object": "chat.completion.chunk",
                "created": created_time,
                "model": self.agent_name,
                "choices": [
                    {
                        "index": 0,
                        "delta": {"content": "[Conversation stopped by user]"},
                        "finish_reason": "stop",
                    }
                ],
            }
            yield f"data: {json.dumps(final_chunk)}\n\n"
            yield "data: [DONE]\n\n"
            logging.info(
                f"Streaming chat completion cancelled for conversation {conversation_id}"
            )
            raise
        except Exception as e:
            logging.error(
                f"Error in streaming chat completions for conversation {conversation_id}: {e}"
            )
            raise
        finally:
            # Always unregister when done
            worker_registry.unregister_conversation(conversation_id)

    async def _execute_chat_completions_stream(self, prompt: ChatCompletions):
        """
        Internal method that does the actual streaming chat completion processing.
        This implementation properly handles tags like <thinking>, <execute>, <answer> etc.
        to ensure backend processing happens while streaming visible content to the frontend.
        """
        from Complexity import (
            calculate_complexity_score,
            should_intervene,
            count_thinking_steps,
            get_planning_phase_prompt,
            get_todo_review_prompt,
            get_answer_review_prompt,
            check_todo_list_exists,
        )

        conversation_id = self.conversation_id
        chunk_id = conversation_id
        created_time = int(time.time())
        c = self.conversation

        # Initialize all parameters like the non-streaming version
        urls = []
        files = []
        new_prompt = ""
        browse_links = True
        tts = False
        websearch = False
        language = "en"
        log_output = True
        log_user_input = True
        disable_commands = False
        running_command = None
        additional_context = ""
        command_overrides = None
        # TTS streaming mode: "off", "audio_only", or "interleaved"
        tts_mode = getattr(prompt, "tts_mode", "off") or "off"

        if prompt.tools:
            command_overrides = prompt.tools
        if "websearch" in self.agent_settings:
            websearch = str(self.agent_settings["websearch"]).lower() == "true"
        if "mode" in self.agent_settings:
            mode = self.agent_settings["mode"]
        else:
            mode = "prompt"
        if "prompt_name" in self.agent_settings:
            prompt_name = self.agent_settings["prompt_name"]
        else:
            prompt_name = "Think About It"
        if "prompt_category" in self.agent_settings:
            prompt_category = self.agent_settings["prompt_category"]
        else:
            prompt_category = "Default"
        if "LANGUAGE" in self.agent_settings:
            language = str(self.agent_settings["LANGUAGE"]).lower()
        prompt_args = {}
        if "prompt_args" in self.agent_settings:
            prompt_args = (
                json.loads(self.agent_settings["prompt_args"])
                if isinstance(self.agent_settings["prompt_args"], str)
                else self.agent_settings["prompt_args"]
            )
        if "context_results" in self.agent_settings:
            context_results = int(self.agent_settings["context_results"])
        else:
            context_results = 5
        if "injected_memories" in self.agent_settings:
            context_results = int(self.agent_settings["injected_memories"])
        if "conversation_results" in self.agent_settings:
            conversation_results = int(self.agent_settings["conversation_results"])
        else:
            conversation_results = 6
        if "command_name" in self.agent_settings:
            command_name = self.agent_settings["command_name"]
        else:
            command_name = ""
        if "command_args" in self.agent_settings:
            try:
                command_args = (
                    json.loads(self.agent_settings["command_args"])
                    if isinstance(self.agent_settings["command_args"], str)
                    else self.agent_settings["command_args"]
                )
            except Exception as e:
                command_args = {}
        else:
            command_args = {}
        if "command_variable" in self.agent_settings:
            command_variable = self.agent_settings["command_variable"]
        else:
            command_variable = "text"
        if "chain_name" in self.agent_settings:
            chain_name = self.agent_settings["chain_name"]
        else:
            chain_name = ""
        if "chain_args" in self.agent_settings:
            chain_args = (
                json.loads(self.agent_settings["chain_args"])
                if isinstance(self.agent_settings["chain_args"], str)
                else self.agent_settings["chain_args"]
            )
        else:
            chain_args = {}
        if "tts_provider" in self.agent_settings:
            tts_provider = str(self.agent_settings["tts_provider"]).lower()
            if tts_provider != "none" and tts_provider != "":
                if "tts" in self.agent_settings:
                    tts = str(self.agent_settings["tts"]).lower() == "true"
        analyze_user_input = False
        if "analyze_user_input" in self.agent_settings:
            analyze_user_input = (
                str(self.agent_settings["analyze_user_input"]).lower() == "true"
            )
        include_sources = False
        if "include_sources" in self.agent_settings:
            include_sources = (
                str(self.agent_settings["include_sources"]).lower() == "true"
            )

        # Extract user message content from the prompt (same as non-streaming)
        parent_activity_id = None
        for message in prompt.messages:
            if "mode" in message:
                if message["mode"] in ["prompt", "command", "chain"]:
                    mode = message["mode"]
            if "log_output" in message:
                log_output = str(message["log_output"]).lower() == "true"
            if "log_user_input" in message:
                log_user_input = str(message["log_user_input"]).lower() == "true"
            if "injected_memories" in message:
                context_results = int(message["injected_memories"])
            if "parent_activity_id" in message:
                parent_activity_id = message["parent_activity_id"]
            if "language" in message:
                language = message["language"]
            if "conversation_results" in message:
                conversation_results = int(message["conversation_results"])
            if "prompt_category" in message:
                prompt_category = message["prompt_category"]
            if "prompt_name" in message:
                prompt_name = message["prompt_name"]
            if "prompt_args" in message:
                prompt_args = (
                    json.loads(message["prompt_args"])
                    if isinstance(message["prompt_args"], str)
                    else message["prompt_args"]
                )
            if "command_name" in message:
                command_name = message["command_name"]
            if "command_args" in message:
                command_args = (
                    json.loads(message["command_args"])
                    if isinstance(message["command_args"], str)
                    else message["command_args"]
                )
            if "command_variable" in message:
                command_variable = message["command_variable"]
            if "chain_name" in message:
                chain_name = message["chain_name"]
            if "chain_args" in message:
                chain_args = (
                    json.loads(message["chain_args"])
                    if isinstance(message["chain_args"], str)
                    else message["chain_args"]
                )
            if "browse_links" in message:
                browse_links = str(message["browse_links"]).lower() == "true"
            if "tts" in message:
                tts = str(message["tts"]).lower() == "true"
            if "websearch" in message:
                websearch = str(message["websearch"]).lower() == "true"
            if "analyze_user_input" in message:
                analyze_user_input = (
                    str(message["analyze_user_input"]).lower() == "true"
                )
            if "context" in message:
                additional_context += "\n" + str(message["context"]).strip()
            if "include_sources" in message:
                include_sources = str(message["include_sources"]).lower() == "true"
            download_headers = {}
            if "download_headers" in message:
                download_headers = (
                    json.loads(message["download_headers"])
                    if isinstance(message["download_headers"], str)
                    else message["download_headers"]
                )
            if "disable_commands" in message:
                disable_commands = str(message["disable_commands"]).lower() == "true"
            if "running_command" in message:
                running_command = message["running_command"]
            if "content" not in message:
                continue
            if isinstance(message["content"], str):
                role = message["role"] if "role" in message else "User"
                if role.lower() == "system":
                    if "/" in message["content"]:
                        new_prompt += f"{message['content']}\n\n"
                if role.lower() == "user":
                    new_prompt += f"{message['content']}\n\n"
            if isinstance(message["content"], list):
                for msg in message["content"]:
                    if "text" in msg:
                        role = message["role"] if "role" in message else "User"
                        if role.lower() == "user":
                            new_prompt += f"{msg['text']}\n\n"
                    # Process file type messages (streaming)
                    await self._process_file_type_message(msg, files)
                    # Iterate over the msg to find _url in one of the keys then use the value of that key unless it has a "url" under it
                    if isinstance(msg, dict):
                        for key, value in msg.items():
                            if "_url" in key:
                                url = str(value["url"] if "url" in value else value)
                                if url.startswith("https://github.com/"):
                                    do_not_pull_repo = [
                                        "/pull/",
                                        "/issues",
                                        "/discussions",
                                        "/actions/",
                                        "/projects",
                                        "/security",
                                        "/releases",
                                        "/commits",
                                        "/branches",
                                        "/tags",
                                        "/stargazers",
                                        "/watchers",
                                        "/network",
                                        "/settings",
                                        "/compare",
                                        "/archive",
                                    ]
                                    if any(x in url for x in do_not_pull_repo):
                                        # If the URL is not a repository, don't pull it
                                        urls.append(url)
                                    else:
                                        # Download the zip for the repo
                                        github_user = (
                                            self.agent_settings["GITHUB_USERNAME"]
                                            if "GITHUB_USERNAME" in self.agent_settings
                                            else None
                                        )
                                        github_token = (
                                            self.agent_settings["GITHUB_TOKEN"]
                                            if "GITHUB_TOKEN" in self.agent_settings
                                            else None
                                        )
                                        github_repo = url.replace(
                                            "https://github.com/", ""
                                        )
                                        github_repo = github_repo.replace(
                                            "https://www.github.com/", ""
                                        )
                                        github_branch = "main"
                                        user = github_repo.split("/")[0]
                                        repo = github_repo.split("/")[1]
                                        if " " in repo:
                                            repo = repo.split(" ")[0]
                                        if "\n" in repo:
                                            repo = repo.split("\n")[0]
                                        # Remove any symbols that would not be in the user, repo, or branch
                                        for symbol in [
                                            " ",
                                            "\n",
                                            "\t",
                                            "\r",
                                            "\\",
                                            "/",
                                            ":",
                                            "*",
                                            "?",
                                            '"',
                                            "<",
                                            ">",
                                        ]:
                                            repo = repo.replace(symbol, "")
                                            user = user.replace(symbol, "")
                                            github_branch = github_branch.replace(
                                                symbol, ""
                                            )
                                        repo_url = f"https://github.com/{user}/{repo}/archive/refs/heads/{github_branch}.zip"
                                        try:
                                            if github_user and github_token:
                                                response = requests.get(
                                                    repo_url,
                                                    auth=(github_user, github_token),
                                                )
                                            else:
                                                response = requests.get(repo_url)
                                        except:
                                            github_branch = "master"
                                            repo_url = f"https://github.com/{user}/{repo}/archive/refs/heads/{github_branch}.zip"
                                            try:
                                                if github_user and github_token:
                                                    response = requests.get(
                                                        repo_url,
                                                        auth=(
                                                            github_user,
                                                            github_token,
                                                        ),
                                                    )
                                                else:
                                                    response = requests.get(repo_url)
                                            except:
                                                pass
                                        if response.status_code == 200:
                                            file_name = (
                                                f"{user}_{repo}_{github_branch}.zip"
                                            )
                                            file_data = response.content
                                            file_path = os.path.normpath(
                                                os.path.join(
                                                    self.agent_workspace,
                                                    conversation_id,
                                                    file_name,
                                                )
                                            )
                                            # Validate path stays within workspace to prevent path traversal
                                            abs_workspace = os.path.abspath(
                                                self.agent_workspace
                                            )
                                            abs_file_path = os.path.abspath(file_path)
                                            if abs_file_path.startswith(
                                                abs_workspace + os.sep
                                            ):
                                                with open(abs_file_path, "wb") as f:
                                                    f.write(file_data)
                                                files.append(
                                                    {
                                                        "file_name": file_name,
                                                        "file_url": f"{self.outputs}/{conversation_id}/{file_name}",
                                                    }
                                                )
                                        else:
                                            urls.append(url)
                                else:
                                    # Not a GitHub URL, check if it's a file or audio
                                    if "file_name" in msg:
                                        file_name = str(msg["file_name"])
                                    else:
                                        file_name = ""
                                    if key != "audio_url":
                                        downloaded_file = (
                                            await self.download_file_to_workspace(
                                                url=url,
                                                file_name=file_name,
                                                download_headers=download_headers,
                                            )
                                        )
                                        if downloaded_file != {}:
                                            files.append(downloaded_file)
                                        else:
                                            urls.append(url)
                                    else:
                                        # If there is an audio_url, it is the user's voice input that needs transcribed before running inference
                                        audio_file_info = (
                                            await self.download_file_to_workspace(
                                                url=url
                                            )
                                        )
                                        if audio_file_info != {}:
                                            full_path = os.path.normpath(
                                                os.path.join(
                                                    self.agent_workspace,
                                                    conversation_id,
                                                    audio_file_info["file_name"],
                                                )
                                            )
                                            if full_path.startswith(
                                                self.agent_workspace
                                            ):
                                                audio_file_path = os.path.join(
                                                    self.agent_workspace,
                                                    conversation_id,
                                                    audio_file_info["file_name"],
                                                )
                                                if os.path.normpath(
                                                    audio_file_path
                                                ).startswith(self.agent_workspace):
                                                    wav_file = os.path.join(
                                                        self.agent_workspace,
                                                        conversation_id,
                                                        f"{uuid.uuid4().hex}.wav",
                                                    )
                                                    AudioSegment.from_file(
                                                        audio_file_path
                                                    ).set_frame_rate(16000).export(
                                                        wav_file, format="wav"
                                                    )
                                                    transcribed_audio = (
                                                        await self.audio_to_text(
                                                            audio_path=wav_file,
                                                        )
                                                    )
                                                    if transcribed_audio:
                                                        new_prompt += transcribed_audio

        # Save the original user prompt before adding file info (for logging)
        original_user_prompt = new_prompt.strip()
        # Add file info to the prompt (for context to the agent)
        for file in files:
            new_prompt += f"\nUploaded file: `{file['file_name']}`."

        # Log user input (log original prompt without file names appended)
        if log_user_input:
            c.log_interaction(role="USER", message=original_user_prompt)

        # Get thinking_id for activity logging
        thinking_id = c.get_thinking_id(agent_name=self.agent_name)

        # Process uploaded files before streaming
        file_contents = []
        current_input_tokens = get_tokens(new_prompt)
        for file in files:
            content = await self.learn_from_file(
                file_url=file["file_url"],
                file_name=file["file_name"],
                user_input=new_prompt,
                collection_id=self.conversation_id,
                thinking_id=thinking_id,
            )
            file_contents.append(content)
        if file_contents:
            file_content = "\n".join(file_contents)
            file_tokens = get_tokens(file_content)
            current_input_tokens = file_tokens + current_input_tokens
        else:
            file_content = ""
            current_input_tokens = self.input_tokens

        # Learn from any URLs that weren't downloaded as files
        await self.learn_from_websites(
            urls=urls,
            summarize_content=False,
        )

        # Handle prompt_args cleanup like non-streaming version
        if "user_input" in prompt_args:
            del prompt_args["user_input"]
        if "prompt_name" in prompt_args:
            prompt_name = prompt_args["prompt_name"]
            del prompt_args["prompt_name"]
        if "prompt_category" in prompt_args:
            prompt_category = prompt_args["prompt_category"]
            del prompt_args["prompt_category"]
        if "websearch" in prompt_args:
            websearch = prompt_args["websearch"]
            del prompt_args["websearch"]
        if "browse_links" in prompt_args:
            browse_links = prompt_args["browse_links"]
            del prompt_args["browse_links"]
        if "tts" in prompt_args:
            tts = prompt_args["voice_response"]
            del prompt_args["tts"]
        if "context_results" in prompt_args:
            context_results = prompt_args["context_results"]
            del prompt_args["context_results"]
        if "conversation_results" in prompt_args:
            conversation_results = prompt_args["conversation_results"]
            del prompt_args["conversation_results"]
        if "analyze_user_input" in prompt_args:
            analyze_user_input = prompt_args["analyze_user_input"]
            del prompt_args["analyze_user_input"]
        if "voice_response" in prompt_args:
            tts = prompt_args["voice_response"]
            del prompt_args["voice_response"]
        if "injected_memories" in prompt_args:
            context_results = prompt_args["injected_memories"]
            del prompt_args["injected_memories"]
        if "shots" in prompt_args:
            del prompt_args["shots"]
        if "data_analysis" in prompt_args:
            del prompt_args["data_analysis"]

        # Send initial streaming chunk
        initial_chunk = {
            "id": chunk_id,
            "object": "chat.completion.chunk",
            "created": created_time,
            "model": self.agent_name,
            "choices": [
                {
                    "index": 0,
                    "delta": {"role": "assistant", "content": ""},
                    "finish_reason": None,
                }
            ],
        }
        yield f"data: {json.dumps(initial_chunk)}\n\n"

        try:
            # Calculate complexity score for inference-time compute scaling
            complexity_score = calculate_complexity_score(
                user_input=new_prompt,
                agent_settings=self.agent_settings,
            )

            # Determine use_smartest based on complexity scoring
            use_smartest = complexity_score.route_to_smartest

            # Build prompt args for processing
            if disable_commands:
                prompt_args["disable_commands"] = True
            if running_command:
                prompt_args["running_command"] = running_command

            # Inject file content into prompt args if within token limits
            if current_input_tokens < self.agent.max_input_tokens:
                if file_content:
                    prompt_args["uploaded_file_data"] = file_content

            # Add additional context if provided
            if additional_context:
                prompt_args["context"] = additional_context

            # Use the streaming inference pipeline to stream tokens in real-time
            # This properly handles thinking/reflection tags and streams answer content
            final_answer = ""

            # Track if we've streamed any answer content progressively
            has_streamed_progressively = False

            # TTS streaming state - stream TTS sentence-by-sentence for real-time audio
            tts_sentence_buffer = ""  # Buffer for accumulating current sentence
            tts_sent_header = False  # Track if we've sent TTS header
            tts_pending_audio = (
                []
            )  # Queue of (sentence, generator) tuples for concurrent TTS

            # Helper to detect sentence boundaries
            def has_complete_sentence(text):
                """Check if text contains a complete sentence ending."""
                import re

                # Match sentence endings: . ! ? followed by space or end, or newlines
                # But ignore abbreviations like "Dr." "Mr." "e.g." etc.
                abbrevs = (
                    r"(?<![A-Z][a-z])\.\s|(?<![a-z]\.[a-z])\.(?:\s|$)|[!?]\s|[!?]$|\n"
                )
                matches = list(re.finditer(abbrevs, text))
                return len(matches) > 0

            def extract_complete_sentences(text):
                """Extract complete sentences and return (sentences, remainder)."""
                import re

                # Find the last sentence boundary
                boundaries = [".", "!", "?", "\n"]
                last_boundary = -1
                for i, char in enumerate(text):
                    if char in boundaries:
                        # Check it's not an abbreviation
                        if char == "." and i > 0 and i < len(text) - 1:
                            # Simple check: if followed by uppercase or space+uppercase, it's a sentence end
                            next_char = text[i + 1] if i + 1 < len(text) else " "
                            if next_char in " \n" or next_char.isupper():
                                last_boundary = i
                        else:
                            last_boundary = i

                if last_boundary >= 0:
                    sentences = text[: last_boundary + 1].strip()
                    remainder = text[last_boundary + 1 :].lstrip()
                    return sentences, remainder
                return "", text

            async for event in self.agent_interactions.run_stream(
                user_input=new_prompt,
                prompt_category=prompt_category,
                prompt_name=prompt_name,
                context_results=context_results,
                conversation_results=conversation_results,
                conversation_name=self.conversation_name,
                conversation_id=self.conversation_id,
                browse_links=browse_links,
                websearch=websearch,
                log_user_input=False,  # Already logged above
                log_output=True,  # Log the final answer to the conversation
                complexity_score=complexity_score,
                use_smartest=use_smartest,
                thinking_id=thinking_id,  # Pass the thinking_id to avoid creating a duplicate
                command_overrides=command_overrides,  # Pass command overrides to enable specific commands
                tts=tts
                or tts_mode != "off",  # Pass TTS flag for filler speech instructions
                **prompt_args,
            ):
                event_type = event.get("type", "")
                content = event.get("content", "")
                is_complete = event.get("complete", False)

                # Stream answer tokens to the frontend via SSE
                if event_type == "answer" and content:
                    if is_complete:
                        # Final answer received - store it
                        final_answer = content
                        # Only send if we haven't been streaming progressively
                        # This handles command results that return in one shot
                        if not has_streamed_progressively:
                            # Stream text chunk (unless audio_only mode)
                            if tts_mode != "audio_only":
                                chunk = {
                                    "id": chunk_id,
                                    "object": "chat.completion.chunk",
                                    "created": created_time,
                                    "model": self.agent_name,
                                    "choices": [
                                        {
                                            "index": 0,
                                            "delta": {"content": content},
                                            "finish_reason": None,
                                        }
                                    ],
                                }
                                yield f"data: {json.dumps(chunk)}\n\n"

                            # Stream TTS for complete answer if TTS mode is enabled
                            if tts_mode in ("audio_only", "interleaved"):
                                try:
                                    import struct

                                    # Clean content for TTS - remove XML tags
                                    tts_content = re.sub(
                                        r"</?(?:thinking|reflection|answer|execute|output|step|reward|count)[^>]*>",
                                        "",
                                        content,
                                        flags=re.IGNORECASE,
                                    )
                                    tts_content = re.sub(
                                        r"(?:thinking|reflection|answer|execute|output|step|reward|count)>",
                                        "",
                                        tts_content,
                                        flags=re.IGNORECASE,
                                    )
                                    # Remove audio HTML tags and URLs that shouldn't be spoken
                                    tts_content = re.sub(
                                        r"<audio[^>]*>.*?</audio>",
                                        "",
                                        tts_content,
                                        flags=re.IGNORECASE | re.DOTALL,
                                    )
                                    tts_content = re.sub(
                                        r"https?://[^\s<>]+",
                                        "",
                                        tts_content,
                                    )
                                    tts_content = re.sub(
                                        r"\s+", " ", tts_content
                                    ).strip()

                                    if not tts_content or len(tts_content) < 2:
                                        continue

                                    raw_buffer = b""
                                    header_sent = False

                                    async for (
                                        audio_chunk
                                    ) in self.agent.text_to_speech_stream(tts_content):
                                        raw_buffer += audio_chunk

                                        # Parse header first (8 bytes)
                                        if not header_sent and len(raw_buffer) >= 8:
                                            header_data = raw_buffer[:8]
                                            raw_buffer = raw_buffer[8:]
                                            header_sent = True
                                            tts_sent_header = True

                                            tts_header_chunk = {
                                                "id": chunk_id,
                                                "object": "audio.header",
                                                "created": created_time,
                                                "model": self.agent_name,
                                                "audio": base64.b64encode(
                                                    header_data
                                                ).decode("utf-8"),
                                            }
                                            yield f"data: {json.dumps(tts_header_chunk)}\n\n"

                                        # Parse data packets: 4-byte size + PCM data
                                        while header_sent and len(raw_buffer) >= 4:
                                            packet_size = struct.unpack(
                                                "<I", raw_buffer[:4]
                                            )[0]
                                            if packet_size == 0:
                                                raw_buffer = raw_buffer[4:]
                                                break
                                            if len(raw_buffer) >= 4 + packet_size:
                                                pcm_data = raw_buffer[
                                                    4 : 4 + packet_size
                                                ]
                                                raw_buffer = raw_buffer[
                                                    4 + packet_size :
                                                ]

                                                # Break large audio chunks into smaller pieces for streaming
                                                # ESP32 has limited buffer size, so send max 4KB at a time
                                                MAX_CHUNK_SIZE = 4096
                                                for offset in range(
                                                    0, len(pcm_data), MAX_CHUNK_SIZE
                                                ):
                                                    chunk_piece = pcm_data[
                                                        offset : offset + MAX_CHUNK_SIZE
                                                    ]
                                                    audio_data_chunk = {
                                                        "id": chunk_id,
                                                        "object": "audio.chunk",
                                                        "created": created_time,
                                                        "model": self.agent_name,
                                                        "audio": base64.b64encode(
                                                            chunk_piece
                                                        ).decode("utf-8"),
                                                    }
                                                    yield f"data: {json.dumps(audio_data_chunk)}\n\n"
                                            else:
                                                break
                                except Exception as e:
                                    logging.warning(f"TTS streaming error: {e}")
                    else:
                        # Progressive answer streaming - send each token
                        has_streamed_progressively = True

                        # Stream text chunk (unless audio_only mode)
                        if tts_mode != "audio_only":
                            chunk = {
                                "id": chunk_id,
                                "object": "chat.completion.chunk",
                                "created": created_time,
                                "model": self.agent_name,
                                "choices": [
                                    {
                                        "index": 0,
                                        "delta": {"content": content},
                                        "finish_reason": None,
                                    }
                                ],
                            }
                            yield f"data: {json.dumps(chunk)}\n\n"

                        # Buffer text for TTS and stream sentence-by-sentence
                        if tts_mode in ("audio_only", "interleaved"):
                            tts_sentence_buffer += content
                            logging.debug(
                                f"[TTS] Buffer: {repr(tts_sentence_buffer[:100])}"
                            )

                            # Check if we have a complete sentence to stream
                            sentences, remainder = extract_complete_sentences(
                                tts_sentence_buffer
                            )
                            if sentences:
                                # Clean TTS sentences - remove any XML tags and fragments
                                # that shouldn't be spoken aloud
                                sentences = re.sub(
                                    r"</?(?:thinking|reflection|answer|execute|output|step|reward|count)[^>]*>",
                                    "",
                                    sentences,
                                    flags=re.IGNORECASE,
                                )
                                # Remove partial tag fragments
                                sentences = re.sub(
                                    r"(?:thinking|reflection|answer|execute|output|step|reward|count)>",
                                    "",
                                    sentences,
                                    flags=re.IGNORECASE,
                                )
                                # Remove audio HTML tags and URLs that shouldn't be spoken
                                sentences = re.sub(
                                    r"<audio[^>]*>.*?</audio>",
                                    "",
                                    sentences,
                                    flags=re.IGNORECASE | re.DOTALL,
                                )
                                sentences = re.sub(
                                    r"https?://[^\s<>]+",
                                    "",
                                    sentences,
                                )
                                # Clean up whitespace left by tag removal
                                sentences = re.sub(r"\s+", " ", sentences).strip()

                                # Skip if after cleaning there's nothing meaningful to speak
                                if not sentences or len(sentences) < 2:
                                    tts_sentence_buffer = remainder
                                    continue

                                logging.info(
                                    f"[TTS] Streaming sentence: {sentences[:80]}..."
                                )
                                tts_sentence_buffer = remainder

                                # Stream TTS for the complete sentence(s)
                                try:
                                    import struct

                                    raw_buffer = b""

                                    async for (
                                        audio_chunk
                                    ) in self.agent.text_to_speech_stream(sentences):
                                        raw_buffer += audio_chunk

                                        # Parse header first (8 bytes) - only once
                                        if not tts_sent_header and len(raw_buffer) >= 8:
                                            header_data = raw_buffer[:8]
                                            raw_buffer = raw_buffer[8:]
                                            tts_sent_header = True

                                            tts_header_chunk = {
                                                "id": chunk_id,
                                                "object": "audio.header",
                                                "created": created_time,
                                                "model": self.agent_name,
                                                "audio": base64.b64encode(
                                                    header_data
                                                ).decode("utf-8"),
                                            }
                                            yield f"data: {json.dumps(tts_header_chunk)}\n\n"

                                        # Parse data packets: 4-byte size + PCM data
                                        while tts_sent_header and len(raw_buffer) >= 4:
                                            packet_size = struct.unpack(
                                                "<I", raw_buffer[:4]
                                            )[0]
                                            if packet_size == 0:
                                                raw_buffer = raw_buffer[4:]
                                                break
                                            if len(raw_buffer) >= 4 + packet_size:
                                                pcm_data = raw_buffer[
                                                    4 : 4 + packet_size
                                                ]
                                                raw_buffer = raw_buffer[
                                                    4 + packet_size :
                                                ]

                                                # Break large audio chunks into smaller pieces for streaming
                                                # ESP32 has limited buffer size, so send max 4KB at a time
                                                MAX_CHUNK_SIZE = 4096
                                                for offset in range(
                                                    0, len(pcm_data), MAX_CHUNK_SIZE
                                                ):
                                                    chunk_piece = pcm_data[
                                                        offset : offset + MAX_CHUNK_SIZE
                                                    ]
                                                    audio_data_chunk = {
                                                        "id": chunk_id,
                                                        "object": "audio.chunk",
                                                        "created": created_time,
                                                        "model": self.agent_name,
                                                        "audio": base64.b64encode(
                                                            chunk_piece
                                                        ).decode("utf-8"),
                                                    }
                                                    yield f"data: {json.dumps(audio_data_chunk)}\n\n"
                                            else:
                                                break
                                except Exception as e:
                                    logging.warning(
                                        f"TTS sentence streaming error: {e}"
                                    )

                # Stream progressive thinking/reflection content
                elif event_type in (
                    "thinking_stream",
                    "reflection_stream",
                    "thinking",
                    "reflection",
                ):
                    # Send as a custom SSE event for progressive activity streaming
                    activity_type = event_type.replace("_stream", "")
                    activity_chunk = {
                        "id": chunk_id,
                        "object": "activity.stream",
                        "type": activity_type,
                        "created": created_time,
                        "content": content,
                        "complete": is_complete,
                    }
                    yield f"data: {json.dumps(activity_chunk)}\n\n"

                # Handle speak events for TTS filler speech during thinking
                # These are brief phrases like "Let me check on that" spoken while processing
                # Audio flows continuously with the main answer - no end marker here
                elif event_type == "speak" and content:
                    if tts_mode in ("audio_only", "interleaved"):
                        logging.info(f"[TTS] Speaking filler: {content}")
                        try:
                            import struct

                            filler_buffer = b""
                            filler_header_sent = False

                            async for audio_chunk in self.agent.text_to_speech_stream(
                                content
                            ):
                                filler_buffer += audio_chunk

                                # Parse header first (8 bytes) - only once per stream
                                if not tts_sent_header and len(filler_buffer) >= 8:
                                    header_data = filler_buffer[:8]
                                    filler_buffer = filler_buffer[8:]
                                    filler_header_sent = True
                                    # Mark main header as sent since format is same
                                    tts_sent_header = True

                                    tts_header_chunk = {
                                        "id": chunk_id,
                                        "object": "audio.header",
                                        "created": created_time,
                                        "model": self.agent_name,
                                        "audio": base64.b64encode(header_data).decode(
                                            "utf-8"
                                        ),
                                    }
                                    yield f"data: {json.dumps(tts_header_chunk)}\n\n"

                                # If we already sent header before, skip this one
                                if (
                                    tts_sent_header
                                    and not filler_header_sent
                                    and len(filler_buffer) >= 8
                                ):
                                    filler_buffer = filler_buffer[8:]
                                    filler_header_sent = True

                                # Parse data packets: 4-byte size + PCM data
                                while (tts_sent_header or filler_header_sent) and len(
                                    filler_buffer
                                ) >= 4:
                                    packet_size = struct.unpack(
                                        "<I", filler_buffer[:4]
                                    )[0]
                                    if packet_size == 0:
                                        filler_buffer = filler_buffer[4:]
                                        break
                                    if len(filler_buffer) >= 4 + packet_size:
                                        pcm_data = filler_buffer[4 : 4 + packet_size]
                                        filler_buffer = filler_buffer[4 + packet_size :]

                                        # Break large audio chunks into smaller pieces for streaming
                                        MAX_CHUNK_SIZE = 4096
                                        for offset in range(
                                            0, len(pcm_data), MAX_CHUNK_SIZE
                                        ):
                                            chunk_piece = pcm_data[
                                                offset : offset + MAX_CHUNK_SIZE
                                            ]
                                            audio_data_chunk = {
                                                "id": chunk_id,
                                                "object": "audio.chunk",
                                                "created": created_time,
                                                "model": self.agent_name,
                                                "audio": base64.b64encode(
                                                    chunk_piece
                                                ).decode("utf-8"),
                                            }
                                            yield f"data: {json.dumps(audio_data_chunk)}\n\n"
                                    else:
                                        break

                            # No audio.end here - let it flow continuously with answer TTS
                            logging.info(f"[TTS] Filler speech sent")
                        except Exception as e:
                            logging.warning(f"TTS speak filler error: {e}")

                # Handle remote command requests - need client-side execution
                elif event_type == "remote_command_request":
                    remote_cmd = content  # content is the remote command dict

                    # Build the SSE chunk - support both legacy format and new client tool format
                    remote_request_chunk = {
                        "id": chunk_id,
                        "object": "remote_command.request",
                        "created": created_time,
                        "model": self.agent_name,
                        "conversation_id": self.conversation_id,
                        "request_id": remote_cmd.get("request_id"),
                    }

                    # Check if this is a client-defined tool call
                    if remote_cmd.get("tool_name"):
                        remote_request_chunk["tool_name"] = remote_cmd.get("tool_name")
                        remote_request_chunk["tool_args"] = remote_cmd.get(
                            "tool_args", {}
                        )
                        # For execute_terminal_command, also include the mapped fields
                        if remote_cmd.get("tool_name") == "execute_terminal_command":
                            remote_request_chunk["command"] = remote_cmd.get(
                                "command", ""
                            )
                            remote_request_chunk["working_directory"] = remote_cmd.get(
                                "working_directory"
                            )
                            remote_request_chunk["terminal_id"] = remote_cmd.get(
                                "terminal_id"
                            )
                    else:
                        # Legacy format
                        remote_request_chunk["command"] = remote_cmd.get("command", "")
                        remote_request_chunk["working_directory"] = remote_cmd.get(
                            "working_directory"
                        )
                        remote_request_chunk["terminal_id"] = remote_cmd.get(
                            "terminal_id"
                        )

                    yield f"data: {json.dumps(remote_request_chunk)}\n\n"

                # Handle remote command pending - stream is ending, waiting for CLI
                elif event_type == "remote_command_pending":
                    pending_chunk = {
                        "id": chunk_id,
                        "object": "remote_command.pending",
                        "created": created_time,
                        "model": self.agent_name,
                        "conversation_id": self.conversation_id,
                    }
                    yield f"data: {json.dumps(pending_chunk)}\n\n"

                elif event_type == "error":
                    error_chunk = {
                        "id": chunk_id,
                        "object": "chat.completion.chunk",
                        "created": created_time,
                        "model": self.agent_name,
                        "choices": [
                            {
                                "index": 0,
                                "delta": {"content": f"Error: {content}"},
                                "finish_reason": "stop",
                            }
                        ],
                    }
                    yield f"data: {json.dumps(error_chunk)}\n\n"

            # Handle conversation rename for new conversations
            if self.conversation_name == "-":
                asyncio.create_task(self.rename_new_conversation(new_prompt))

            # Generate TTS for any remaining buffered text (partial sentences at end)
            if (
                tts_mode in ("audio_only", "interleaved")
                and tts_sentence_buffer.strip()
            ):
                try:
                    import struct

                    # Clean remaining TTS buffer - remove XML tags
                    final_tts_text = re.sub(
                        r"</?(?:thinking|reflection|answer|execute|output|step|reward|count)[^>]*>",
                        "",
                        tts_sentence_buffer,
                        flags=re.IGNORECASE,
                    )
                    final_tts_text = re.sub(
                        r"(?:thinking|reflection|answer|execute|output|step|reward|count)>",
                        "",
                        final_tts_text,
                        flags=re.IGNORECASE,
                    )
                    final_tts_text = re.sub(r"\s+", " ", final_tts_text).strip()

                    # Skip if nothing meaningful remains
                    if final_tts_text and len(final_tts_text) >= 2:
                        # Buffer to accumulate incoming bytes and parse properly
                        raw_buffer = b""

                        async for audio_chunk in self.agent.text_to_speech_stream(
                            final_tts_text
                        ):
                            raw_buffer += audio_chunk

                            # Parse header first (8 bytes) - only if not already sent
                            if not tts_sent_header and len(raw_buffer) >= 8:
                                header_data = raw_buffer[:8]
                                raw_buffer = raw_buffer[8:]
                                tts_sent_header = True

                                tts_header_chunk = {
                                    "id": chunk_id,
                                    "object": "audio.header",
                                    "created": created_time,
                                    "model": self.agent_name,
                                    "audio": base64.b64encode(header_data).decode(
                                        "utf-8"
                                    ),
                                }
                                yield f"data: {json.dumps(tts_header_chunk)}\n\n"

                            # Parse data packets: 4-byte size + PCM data
                            while tts_sent_header and len(raw_buffer) >= 4:
                                # Read packet size
                                packet_size = struct.unpack("<I", raw_buffer[:4])[0]

                                # Check for end marker
                                if packet_size == 0:
                                    raw_buffer = raw_buffer[4:]
                                    break

                                # Check if we have the full packet
                                if len(raw_buffer) >= 4 + packet_size:
                                    pcm_data = raw_buffer[4 : 4 + packet_size]
                                    raw_buffer = raw_buffer[4 + packet_size :]

                                    # Break large audio chunks into smaller pieces for streaming
                                    # ESP32 has limited buffer size, so send max 4KB at a time
                                    MAX_CHUNK_SIZE = 4096
                                    for offset in range(
                                        0, len(pcm_data), MAX_CHUNK_SIZE
                                    ):
                                        chunk_piece = pcm_data[
                                            offset : offset + MAX_CHUNK_SIZE
                                        ]
                                        audio_data_chunk = {
                                            "id": chunk_id,
                                            "object": "audio.chunk",
                                            "created": created_time,
                                            "model": self.agent_name,
                                            "audio": base64.b64encode(
                                                chunk_piece
                                            ).decode("utf-8"),
                                        }
                                        yield f"data: {json.dumps(audio_data_chunk)}\n\n"
                                else:
                                    # Not enough data yet, wait for more
                                    break

                except Exception as e:
                    logging.warning(f"TTS streaming error: {e}")

            # Send TTS end marker if we sent any TTS data
            if tts_sent_header:
                tts_end_chunk = {
                    "id": chunk_id,
                    "object": "audio.end",
                    "created": created_time,
                    "model": self.agent_name,
                }
                yield f"data: {json.dumps(tts_end_chunk)}\n\n"

        except Exception as e:
            logging.error(f"Streaming error: {str(e)}")
            import traceback

            logging.error(traceback.format_exc())
            # Send generic error chunk (do NOT expose exception details!)
            error_chunk = {
                "id": chunk_id,
                "object": "chat.completion.chunk",
                "created": created_time,
                "model": self.agent_name,
                "choices": [
                    {
                        "index": 0,
                        "delta": {"content": "Error: An internal error has occurred."},
                        "finish_reason": "stop",
                    }
                ],
            }
            yield f"data: {json.dumps(error_chunk)}\n\n"

        # Send final chunk to indicate completion
        final_chunk = {
            "id": chunk_id,
            "object": "chat.completion.chunk",
            "created": created_time,
            "model": self.agent_name,
            "choices": [
                {
                    "index": 0,
                    "delta": {},
                    "finish_reason": "stop",
                }
            ],
        }
        yield f"data: {json.dumps(final_chunk)}\n\n"
        yield "data: [DONE]\n\n"

    async def batch_inference(
        self,
        user_inputs: List[str] = [],
        prompt_category: str = "Default",
        prompt_name: str = "Ask Questions",
        images: list = [],
        injected_memories: int = 100,
        batch_size: int = 10,
        browse_links: bool = False,
        voice_response: bool = False,
        log_user_input: bool = False,
        **kwargs,
    ):
        i = 0
        tasks = []
        responses = []
        if user_inputs == []:
            return []
        for user_input in user_inputs:
            i += 1
            if i % batch_size == 0:
                responses += await asyncio.gather(**tasks)
                tasks = []
            task = asyncio.create_task(
                await self.inference(
                    user_input=user_input,
                    prompt_category=prompt_category,
                    prompt_name=prompt_name,
                    images=images,
                    injected_memories=injected_memories,
                    browse_links=browse_links,
                    voice_response=voice_response,
                    log_user_input=log_user_input,
                    **kwargs,
                )
            )
            tasks.append(task)
        responses += await asyncio.gather(**tasks)
        return responses

    async def dpo(
        self,
        question: str = "",
        injected_memories: int = 100,
    ):
        context_async = self.memories(
            user_input=question,
            limit_per_collection=injected_memories,
        )
        chosen_async = self.inference(
            user_input=question,
            prompt_category="Default",
            prompt_name="Answer Question with Memory",
            injected_memories=injected_memories,
            log_user_input=False,
            log_output=False,
        )
        rejected_async = self.inference(
            user_input=question,
            prompt_category="Default",
            prompt_name="Wrong Answers Only",
            log_user_input=False,
            log_output=False,
        )
        chosen = await chosen_async
        rejected = await rejected_async
        context = await context_async
        prompt = f"### Context\n{context}\n### Question\n{question}"
        return prompt, chosen, rejected

    # Creates a synthetic dataset from memories in sharegpt format
    async def create_dataset_from_memories(self, batch_size: int = 10):
        self.agent_settings["training"] = True
        self.agent_interactions.agent.update_agent_config(
            new_config=self.agent_settings, config_key="settings"
        )
        memories = []
        questions = []
        dataset_name = f"{datetime.now().strftime('%Y-%m-%d-%H-%M-%S')}-DPO-Dataset"
        collections = await self.agent_interactions.agent_memory.get_collections()
        for collection in collections:
            self.collection_name = collection
            memories += (
                await self.agent_interactions.agent_memory.export_collection_to_json()
            )
        memories = [memory["text"] for memory in memories]
        # Get a list of questions about each memory
        question_list = self.batch_inference(
            user_inputs=memories,
            batch_size=batch_size,
            prompt_category="Default",
            prompt_name="Ask Questions",
        )
        for question in question_list:
            # Convert the response to a list of questions
            question = question.split("\n")
            question = [
                item.lstrip("0123456789.*- ") for item in question if item.lstrip()
            ]
            question = [item for item in question if item]
            question = [item.lstrip("0123456789.*- ") for item in question]
            questions += question
        prompts = []
        good_answers = []
        bad_answers = []
        for question in questions:
            prompt, chosen, rejected = await self.dpo(
                question=question,
                injected_memories=100,
            )
            prompts.append(prompt)
            good_answers.append(
                [
                    {"content": prompt, "role": "user"},
                    {"content": chosen, "role": "assistant"},
                ]
            )
            bad_answers.append(
                [
                    {"content": prompt, "role": "user"},
                    {"content": rejected, "role": "assistant"},
                ]
            )
        dpo_dataset = {
            "prompt": questions,
            "chosen": good_answers,
            "rejected": bad_answers,
        }
        # Save messages to a json file to be used as a dataset
        dataset_dir = os.path.join(self.agent_workspace, "datasets")

        os.makedirs(dataset_dir, exist_ok=True)
        dataset_name = "".join(
            [c for c in dataset_name if c.isalpha() or c.isdigit() or c == " "]
        )
        dataset_filename = f"{dataset_name}.json"
        full_path = os.path.normpath(
            os.path.join(self.agent_workspace, dataset_filename)
        )
        if not full_path.startswith(self.agent_workspace):
            raise Exception("Path given not allowed")
        with open(os.path.join(dataset_dir, dataset_filename), "w") as f:
            f.write(json.dumps(dpo_dataset))
        self.agent_settings["training"] = False
        self.agent_interactions.agent.update_agent_config(
            new_config=self.agent_settings, config_key="settings"
        )
        return dpo_dataset

    def _generate_detailed_schema(self, model: Type[BaseModel], depth: int = 0) -> str:
        """
        Recursively generates a detailed schema representation of a Pydantic model,
        including nested models and complex types.
        """
        fields = get_type_hints(model)
        field_descriptions = []
        indent = "  " * depth
        for field, field_type in fields.items():
            description = f"{indent}{field}: "
            origin_type = get_origin(field_type)
            if origin_type is None:
                origin_type = field_type
            if inspect.isclass(origin_type) and issubclass(origin_type, BaseModel):
                description += f"Nested Model:\n{self._generate_detailed_schema(origin_type, depth + 1)}"
            elif origin_type == list:
                list_type = get_args(field_type)[0]
                if inspect.isclass(list_type) and issubclass(list_type, BaseModel):
                    description += f"List of Nested Model:\n{self._generate_detailed_schema(list_type, depth + 1)}"
                elif get_origin(list_type) == Union:
                    union_types = get_args(list_type)
                    description += f"List of Union:\n"
                    for union_type in union_types:
                        if inspect.isclass(union_type) and issubclass(
                            union_type, BaseModel
                        ):
                            description += f"{indent}  - Nested Model:\n{self._generate_detailed_schema(union_type, depth + 2)}"
                        else:
                            description += (
                                f"{indent}  - {self._get_type_name(union_type)}\n"
                            )
                else:
                    description += f"List[{self._get_type_name(list_type)}]"
            elif origin_type == dict:
                key_type, value_type = get_args(field_type)
                description += f"Dict[{self._get_type_name(key_type)}, {self._get_type_name(value_type)}]"
            elif origin_type == Union:
                union_types = get_args(field_type)

                for union_type in union_types:
                    if inspect.isclass(union_type) and issubclass(
                        union_type, BaseModel
                    ):
                        description += f"{indent}  - Nested Model:\n{self._generate_detailed_schema(union_type, depth + 2)}"
                    else:
                        type_name = self._get_type_name(union_type)
                        if type_name != "NoneType":
                            description += f"{self._get_type_name(union_type)}\n"
            elif inspect.isclass(origin_type) and issubclass(origin_type, Enum):
                enum_values = ", ".join([f"{e.name} = {e.value}" for e in origin_type])
                description += f"{origin_type.__name__} (Enum values: {enum_values})"
            else:
                description += self._get_type_name(origin_type)
            field_descriptions.append(description)
        return "\n".join(field_descriptions)

    def _get_type_name(self, type_):
        """Helper method to get the name of a type, handling some special cases."""
        if hasattr(type_, "__name__"):
            return type_.__name__
        return str(type_).replace("typing.", "")

    async def convert_to_model(
        self,
        input_string: str,
        model: Type[BaseModel],
        max_failures: int = 3,
        response_type: str = None,
        **kwargs,
    ):
        """
        Converts a string to a Pydantic model using an AGiXT agent.

        Args:
        input_string (str): The string to convert to a model.
        model (Type[BaseModel]): The Pydantic model to convert the string to.
        max_failures (int): The maximum number of times to retry the conversion if it fails.
        response_type (str): The type of response to return. Either 'json' or None. None will return the model.
        **kwargs: Additional arguments to pass to the AGiXT agent as prompt arguments.
        """
        input_string = str(input_string)
        schema = self._generate_detailed_schema(model)
        if "user_input" in kwargs:
            del kwargs["user_input"]
        if "schema" in kwargs:
            del kwargs["schema"]
        response = await self.inference(
            user_input=input_string,
            schema=schema,
            prompt_category="Default",
            prompt_name="Convert to Pydantic Model",
            log_user_input=False,
            log_output=False,
        )
        if "```json" in response:
            response = response.split("```json")[1].split("```")[0].strip()
        elif "```" in response:
            response = response.split("```")[1].strip()
        try:
            response = json.loads(response)
            if response_type == "json":
                return response
            else:
                return model(**response)
        except Exception as e:
            if "failures" in kwargs:
                failures = int(kwargs["failures"]) + 1
                if failures > max_failures:
                    logging.error(
                        f"Error: {e} . Failed to convert the response to the model after 3 attempts. Response: {response}"
                    )
                    return (
                        response
                        if response
                        else "Failed to convert the response to the model."
                    )
            else:
                failures = 1
            logging.warning(
                f"Error: {e} . Failed to convert the response to the model, trying again. {failures}/3 failures. Response: {response}"
            )
            return await self.convert_to_model(
                input_string=input_string,
                model=model,
                max_failures=max_failures,
                response_type=response_type,
                failures=failures,
            )

    def get_agent_workspace_markdown(self):
        def generate_markdown_structure(folder_path, indent=0):
            if not os.path.isdir(folder_path):
                return ""
            markdown_output = ""
            items = sorted(os.listdir(folder_path))
            for item in items:
                item_path = os.path.join(folder_path, item)
                if os.path.isdir(item_path):
                    markdown_output += f"{'  ' * indent}* **{item}/**\n"
                    markdown_output += generate_markdown_structure(
                        item_path, indent + 1
                    )
                else:
                    markdown_output += f"{'  ' * indent}* {item}\n"
            return markdown_output

        return generate_markdown_structure(folder_path=self.conversation_workspace)

    def get_agent_workspace_list(self):
        files_with_paths = []
        for root, dirs, files in os.walk(self.conversation_workspace):
            for file in files:
                files_with_paths.append(os.path.join(root, file))
        files_with_paths.sort()
        return files_with_paths

    async def analyze_user_input(self, user_input: str):
        code_execution = ""
        thinking_id = self.conversation.get_thinking_id(agent_name=self.agent_name)
        self.conversation.log_interaction(
            role=self.agent_name,
            message=f"[SUBACTIVITY][{thinking_id}] Analyzing.",
        )
        analyze_input = await self.inference(
            user_input=user_input,
            prompt_category="Default",
            prompt_name="Analyze Input",
            log_user_input=False,
            log_output=False,
            browse_links=False,
            websearch=False,
            websearch_depth=0,
            voice_response=False,
        )
        analyzed_input = {}
        if "```json" not in analyze_input and "```" in analyze_input:
            analyze_input = analyze_input.replace("```", "```json", 1)
        if "```json" in analyze_input:
            analyze_input = analyze_input.split("```json")[1].split("```")[0]
        try:
            analyzed_input = json.loads(analyze_input)
            if "math" in analyzed_input:
                if str(analyzed_input["math"]).lower() != "true":
                    return ""
            else:
                return ""
        except:
            return ""
        if analyzed_input == {}:
            return ""
        code_interpreter = await self.inference(
            user_input=user_input,
            prompt_category="Default",
            prompt_name="Write Code",
            log_user_input=False,
            log_output=False,
            browse_links=False,
            websearch=False,
            websearch_depth=0,
            voice_response=False,
        )
        if "```" not in code_interpreter:
            code_interpreter = f"```python\n{code_interpreter}\n```"
        if "```python" in code_interpreter:
            code_interpreter = code_interpreter.split("```python")[1].split("```")[0]
            if "```python" in code_interpreter:
                code_interpreter = code_interpreter.split("```python")[1].split("```")[
                    0
                ]
            code_verification = await self.inference(
                user_input=user_input,
                prompt_category="Default",
                prompt_name="Verify Code",
                code=code_interpreter,
                log_user_input=False,
                log_output=False,
                browse_links=False,
                websearch=False,
                websearch_depth=0,
                voice_response=False,
            )
            if "```" not in code_verification:
                code_verification = f"```python\n{code_verification}\n```"
            if "```python" in code_verification:
                code_verification = code_verification.split("```python")[1].split(
                    "```"
                )[0]
                if "```python" in code_verification:
                    code_verification = code_verification.split("```python")[1].split(
                        "```"
                    )[0]
                try:
                    last_line = code_verification.split("\n")[-1]
                    if last_line == "\n" or last_line == "":
                        last_line = code_verification.split("\n")[-2]
                    if last_line == "\n" or last_line == "":
                        last_line = code_verification.split("\n")[-3]
                    if not last_line.startswith("print("):
                        old_last_line = last_line
                        new_last_line = f"print({last_line})"
                        code_verification = code_verification.rsplit(old_last_line, 1)[
                            0
                        ]
                        code_verification += new_last_line
                except Exception as e:
                    logging.error(f"Error adding print statement: {e}")
                try:
                    code_execution = await self.execute_command(
                        command_name="Execute Python Code",
                        command_args={"code": code_verification, "text": ""},
                    )
                except Exception as e:
                    fixed_code = await self.inference(
                        user_input=user_input,
                        prompt_category="Default",
                        prompt_name="Fix Verified Code",
                        code=code_verification,
                        code_error=str(e),
                        log_user_input=False,
                        log_output=False,
                        browse_links=False,
                        websearch=False,
                        websearch_depth=0,
                        voice_response=False,
                    )
                    if "```" not in fixed_code:
                        fixed_code = f"```python\n{fixed_code}\n```"
                    code_verification = fixed_code
                    if "```python" in code_verification:
                        code_verification = code_verification.split("```python")[
                            1
                        ].split("```")[0]
                        if "```python" in code_verification:
                            code_verification = code_verification.split("```python")[
                                1
                            ].split("```")[0]
                        try:
                            code_execution = await self.execute_command(
                                command_name="Execute Python Code",
                                command_args={"code": code_verification, "text": ""},
                            )
                        except Exception as e:
                            code_execution = ""
                            logging.error(f"Error executing code: {e}")
        if code_execution != "":
            return f"Executed the following code expressed to assist the user:\n```python\n{code_verification}\n```\n**REFERENCE THE RESULTS, NOT THE CODE TO THE USER WHEN RESPONDING! THE RESULTS FROM RUNNING THE CODE ARE AS FOLLOWS:**\n{code_execution}"
        return ""

    async def fix_and_execute_code(
        self,
        user_input: str,
        code: str,
        code_error: str,
        file_content: str,
        file_preview: str = "",
        import_file: str = "",
        multifile: bool = False,
        failures: int = 0,
        max_failures: int = 5,
    ):
        code_failed = False
        fixed_code = await self.inference(
            user_input=user_input,
            prompt_category="Default",
            prompt_name="Fix Code",
            file_preview=file_preview,
            import_file=import_file,
            code=code,
            code_error=str(code_error),
            log_user_input=False,
            log_output=False,
            browse_links=False,
            websearch=False,
            websearch_depth=0,
            voice_response=False,
        )
        if "```python" in fixed_code:
            fixed_code = fixed_code.split("```python")[1].split("```")[0]
        if "```python" in fixed_code:
            fixed_code = fixed_code.split("```python")[1].split("```")[0]
        code_verification = await self.inference(
            user_input=user_input,
            prompt_category="Default",
            prompt_name=(
                "Verify Code Interpreter Multifile"
                if multifile
                else "Verify Code Interpreter"
            ),
            import_file=import_file,
            file_preview=file_preview,
            code=fixed_code,
            log_user_input=False,
            log_output=False,
            browse_links=False,
            websearch=False,
            websearch_depth=0,
            voice_response=False,
        )
        try:
            code_execution = await self.execute_command(
                command_name="Execute Python Code",
                command_args={"code": code_verification, "text": file_content},
            )
        except Exception:
            code_failed = True
        if code_execution.startswith("Error"):
            code_failed = True
        if code_failed:
            failures += 1
            if failures >= max_failures:
                return code_verification
            return await self.fix_and_execute_code(
                user_input=user_input,
                code=code_verification,
                code_error=str(code_execution),
                file_content=file_content,
                file_preview=file_preview,
                import_file=import_file,
                failures=failures,
                max_failures=max_failures,
            )

    async def analyze_data(
        self,
        user_input: str,
        file_content=None,
        file_name="",
        max_failures: int = 5,
    ):
        file_names = []
        import_files = ""
        file_preview = ""
        file_path = self.conversation_workspace
        thinking_id = self.conversation.get_thinking_id(agent_name=self.agent_name)
        if "```csv" in user_input and file_name == "":
            file_name = f"{uuid.uuid4().hex}.csv"
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[SUBACTIVITY][{thinking_id}] Saving CSV data to file `{file_name}`.",
            )
            file_content = user_input.split("```csv")[1].split("```")[0]
            file_path = os.path.join(self.conversation_workspace, file_name)
            with open(file_path, "w") as f:
                f.write(file_content)
            file_names.append(file_name)
        files = self.get_agent_workspace_list()
        if len(files) != 0:
            csv_files = []
            for file in files:
                if str(file).endswith(".csv"):
                    csv_files.append(file)
            if len(csv_files) != 0:
                for file in csv_files:
                    file_names.append(file)
        if len(file_names) == 0:
            return await self.analyze_user_input(user_input=user_input)
        # Iterate over files and use regex to see if the file name is in the response
        if len(file_names) == 1:
            file_name = file_names[0]
            if self.conversation_workspace not in file_name:
                file_path = os.path.join(self.conversation_workspace, file_name)
            else:
                file_path = file_name
            file_content = open(file_path, "r").read()
            lines = file_content.split("\n")
            if len(lines) < 20:
                file_preview = f"`{file_path}`\n```csv\n{file_content}\n```"
            else:
                limited_content = "\n".join(lines[0:20])
                file_preview = f"`{file_path}`\n```csv\n{limited_content}\n```"
        if len(file_names) > 1:
            # Found multiple files, do things a little differently.
            previews = []
            for file in file_names:
                if self.conversation_workspace not in file:
                    file_path = os.path.join(self.conversation_workspace, file)
                else:
                    file_path = file
                if import_files == "":
                    import_files = f"`{file_path}`"
                else:
                    import_files += f", `{file_path}`"
                file_content = open(file_path, "r").read()
                lines = file_content.split("\n")
                if len(lines) < 20:
                    previews.append(f"`{file_path}`\n```csv\n{file_content}\n```")
                else:
                    limited_content = "\n".join(lines[0:20])
                    previews.append(f"`{file_path}`\n```csv\n{limited_content}\n```")
            file_preview = "\n".join(previews)
        code_interpreter = await self.inference(
            user_input=user_input,
            prompt_category="Default",
            prompt_name=(
                "Code Interpreter Multifile"
                if len(file_names) > 1
                else "Code Interpreter"
            ),
            import_file=import_files if len(file_names) > 1 else file_path,
            file_preview=file_preview,
            log_user_input=False,
            log_output=False,
            browse_links=False,
            websearch=False,
            websearch_depth=0,
            voice_response=False,
        )
        if "```python" in code_interpreter:
            code_interpreter = code_interpreter.split("```python")[1].split("```")[0]
        if "```python" in code_interpreter:
            code_interpreter = code_interpreter.split("```python")[1].split("```")[0]
        # Step 5 - Verify the code is good before executing it.
        import_file = import_files if len(file_names) > 1 else file_path
        code_verification = await self.inference(
            user_input=user_input,
            prompt_category="Default",
            prompt_name=(
                "Verify Code Interpreter Multifile"
                if len(file_names) > 1
                else "Verify Code Interpreter"
            ),
            import_file=import_file,
            file_preview=file_preview,
            code=code_interpreter,
            log_user_input=False,
            log_output=False,
            browse_links=False,
            websearch=False,
            websearch_depth=0,
            voice_response=False,
        )
        # Split out the python code
        if "```python" in code_verification:
            code_verification = code_verification.split("```python")[1].split("```")[0]
        if "```python" in code_verification:
            code_verification = code_verification.split("```python")[1].split("```")[0]
        # Step 6 - Execute the code
        code_failed = False
        try:
            code_execution = await self.execute_command(
                command_name="Execute Python Code",
                command_args={"code": code_verification, "text": file_content},
            )
        except Exception as e:
            code_failed = True
        if not code_execution:
            code_failed = True
        else:
            if code_execution.startswith("Error"):
                code_failed = True
        # Step 7 - If the code failed to run without error, attempt fix the code and try again {max_failures} times
        if code_failed:
            code_execution = await self.fix_and_execute_code(
                user_input=user_input,
                code=code_verification,
                code_error=str(code_execution),
                file_content=file_content,
                file_preview=file_preview,
                import_file=import_file,
                max_failures=max_failures,
            )
        if code_execution.startswith("Error"):
            self.conversation.log_interaction(
                role=self.agent_name,
                message=f"[SUBACTIVITY][{thinking_id}][ERROR] Unable to complete data analysis.",
            )
            return f"Data analysis failed after {max_failures} attempts. Advise the user that there may be an issue with the data and to try again in a new conversation."
        return f"**REFERENCE ALL OF THE FOLLOWING OUTPUT FROM DATA ANALYSIS RESULTS ON {import_files if len(file_names) > 1 else file_path} INCLUDING ALL VISUALIZATIONS IN MARKDOWN FORMAT TO THE USER. REFERENCE EXACT LINKS TO IMAGES OR FILES IF PRESENT HERE! Do not rename files!**\n\n{code_execution}"
